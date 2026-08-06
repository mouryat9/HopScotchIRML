"""
Hopscotch IRML - chat backend with LLM + RAG.

Flow:
- /session          -> create session
- /chat/history     -> get full chat history
- /chat/send        -> main chat endpoint (non-streaming)
- /chat/send_stream -> streaming endpoint

Chat behaviour:
- Worldview is selected via Step 1 dropdown (/worldview/set).
- All messages go to LLM (Ollama) with:
    - user's worldview band
    - retrieved IRML snippets (RAG, FAISS)
    - step-specific guidance
- Assistant responds in a short, friendly, tutor-style format.
"""

from __future__ import annotations

from typing import List, Dict, Optional, Literal, Any
from pathlib import Path
from datetime import datetime
import uuid
import re
import json
import logging

import requests
from fastapi import FastAPI, HTTPException, Body, Query, Depends, Request, UploadFile, File, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, Response
from pydantic import BaseModel, Field
from jinja2 import Template
from weasyprint import HTML

import os
import resend

from auth import (
    hash_password, verify_password, create_access_token, get_current_user,
    create_password_reset_token, decode_token, require_admin,
)
from database import (
    ensure_indexes, find_user_by_email, find_user_by_username,
    find_user_by_id,
    create_user, create_classroom_student, update_user_password,
    find_session, create_session_doc, update_session,
    get_sessions_for_user,
    get_latest_session_for_user,
    get_session_summaries_for_user,
    create_class_doc, find_class_by_code, get_classes_for_teacher,
    find_class_by_id, update_class_settings, get_class_settings,
    get_students_in_class, get_all_student_sessions_for_teacher,
    # Admin / login tracking
    record_login, get_recent_logins, get_login_locations,
    get_logins_for_user, get_all_users, update_user_fields,
    delete_user_by_id, get_user_counts_by_role, get_signups_over_time,
    get_total_sessions_count, get_total_classes_count,
    get_active_users_last_n_days, get_step_completion_across_all,
    record_admin_action, get_admin_audit_log,
    # Admin: classes, sessions, user detail, geo stats
    get_all_classes, delete_class_by_id,
    get_class_detail, update_class_fields, set_class_students_password,
    get_session_stats, delete_session_by_sid, delete_sessions_bulk,
    get_all_sessions, get_session_full,
    get_user_detail,
    get_login_stats_by_country, get_login_stats_by_region,
    # Glossary
    get_all_glossary_terms, count_glossary_terms, create_glossary_term,
    update_glossary_term, delete_glossary_term, seed_glossary_if_empty,
    set_glossary_translation, get_glossary_ids_missing, get_glossary_term_by_id,
    # Step resources
    get_step_resources, get_step_resources_all, upsert_step_resource,
    seed_step_resources_if_empty,
)

# -------------------------------------------------
# Paths
# -------------------------------------------------
ROOT = Path(__file__).parent
PATHS_PATH = ROOT / "server" / "config" / "paths" / "research_paths.json"
PATHS_ES_PATH = ROOT / "server" / "config" / "paths" / "research_paths.es.json"
PATHS_ZH_PATH = ROOT / "server" / "config" / "paths" / "research_paths.zh.json"
PATHS_OVERLAY_FILES = {"es": PATHS_ES_PATH, "zh": PATHS_ZH_PATH}
TEMPLATE_DIR = ROOT / "server" / "templates"

# -------------------------------------------------
# Email / Resend configuration (password reset)
# -------------------------------------------------
RESEND_API_KEY = os.environ.get("RESEND_API_KEY", "")
EMAIL_FROM = "noreply@hopscotch4all.com"
FRONTEND_URL = os.environ.get("HOPSCOTCH_FRONTEND_URL", "https://hopscotchai.us")

DOCS_DIR = ROOT / "server" / "resources"
INDEX_DIR = ROOT / "server" / "index"
INDEX_DIR.mkdir(parents=True, exist_ok=True)
INDEX_PATH = INDEX_DIR / "faiss.index"
META_PATH = INDEX_DIR / "chunks.json"  # keep chunk texts

# -------------------------------------------------
# Optional RAG deps (safe import)
# -------------------------------------------------
RAG_AVAILABLE = True
try:
    import faiss  # type: ignore
    from sentence_transformers import SentenceTransformer  # type: ignore
    from pypdf import PdfReader  # type: ignore
except Exception:
    RAG_AVAILABLE = False
    faiss = None
    SentenceTransformer = None
    PdfReader = None

# Optional pdfminer fallback (used only if installed)
try:
    from pdfminer_high_level import extract_text as pdfminer_extract_text  # type: ignore
except Exception:
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract_text  # type: ignore
    except Exception:
        pdfminer_extract_text = None

EMBED_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
EMBED_DIM = 384  # for the model above

logger = logging.getLogger("uvicorn.error")

# -------------------------------------------------
# Ollama / LLM config
# -------------------------------------------------
# LLM backend — supports vLLM (primary) with Ollama fallback
LLM_BACKEND = os.environ.get("LLM_BACKEND", "ollama")  # "vllm" or "ollama"
VLLM_URL = os.environ.get("VLLM_URL", "http://127.0.0.1:8000/v1/chat/completions")
VLLM_API_KEY = os.environ.get("VLLM_API_KEY", "")
OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://127.0.0.1:11434/api/chat")
OLLAMA_BASE = os.environ.get("OLLAMA_BASE", "http://127.0.0.1:11434")
LLM_MODEL = os.environ.get("LLM_MODEL", "qwen2.5:14b")
VLLM_MODEL = os.environ.get("VLLM_MODEL", "Qwen/Qwen2.5-14B-Instruct")
LLM_TEMP = float(os.environ.get("LLM_TEMP", "0.4"))

# Harmful-content moderation (Llama Guard 3 via Ollama). Runs locally/free.
# Set MODERATION_ENABLED=0 to disable, or point MODERATION_MODEL at another guard model.
MODERATION_ENABLED = os.environ.get("MODERATION_ENABLED", "1") not in ("0", "false", "False", "")
MODERATION_MODEL = os.environ.get("MODERATION_MODEL", "llama-guard3:1b")

import time as _time_mod
_SERVER_START_TIME = _time_mod.time()

# runtime globals for RAG
_embedder = None
_faiss_index = None
_chunks: List[Dict[str, Any]] = []  # [{"id": int, "text": str, "source": str}]
_raw_docs_cache: Optional[List[Dict[str, str]]] = None  # for keyword fallback

# runtime global for path config
_paths_config: Dict[str, Any] = {}


def load_paths_config() -> Dict[str, Any]:
    """Load research_paths.json once and cache it."""
    global _paths_config
    if _paths_config:
        return _paths_config
    try:
        with open(PATHS_PATH, "r", encoding="utf-8") as f:
            _paths_config = json.load(f)
    except FileNotFoundError:
        logger.warning("Paths config not found at %s", PATHS_PATH)
        _paths_config = {}
    except json.JSONDecodeError as e:
        logger.warning("Paths config JSON invalid: %s", e)
        _paths_config = {}
    return _paths_config


_paths_overlays: Dict[str, Dict[str, Any]] = {}

def load_paths_overlay(lang: str) -> Dict[str, Any]:
    """Localized overlay for student-visible step-config strings (optional file)."""
    if lang not in _paths_overlays:
        try:
            with open(PATHS_OVERLAY_FILES[lang], "r", encoding="utf-8") as f:
                _paths_overlays[lang] = json.load(f)
        except Exception:
            _paths_overlays[lang] = {}
    return _paths_overlays[lang]


def _localize_step_cfg(path_name: str, step_key: str, step_cfg: dict, lang: str) -> dict:
    """Merge a language overlay over a step config. Overlay keys replace the
    English ones wholesale (arrays keep the same ids); anything untranslated
    falls through to English. llm_guidance is never overlaid."""
    if lang not in PATHS_OVERLAY_FILES or not step_cfg:
        return step_cfg
    ov = (load_paths_overlay(lang).get("paths", {}).get(path_name, {})
          .get("steps", {}).get(step_key) or {})
    return {**step_cfg, **ov} if ov else step_cfg



# -------------------------------------------------
# ============================================================
# Models
# ============================================================
class ChatTurn(BaseModel):
    role: Literal["user", "assistant"]
    content: str
    step: Optional[int] = None


class SessionData(BaseModel):
    id: str
    created_at: str
    chat: List[ChatTurn] = Field(default_factory=list)

    worldview_band: Optional[str] = None
    worldview_label: Optional[str] = None

    # arbitrary notes/data per step (1-9)
    step_notes: Dict[str, Any] = Field(default_factory=dict)

    # resolved research path ("quantitative" | "qualitative" | "mixed")
    resolved_path: Optional[str] = None
    # for mixed-methods students: which methodology they chose at Step 4
    chosen_methodology: Optional[str] = None

    # current step the student is working on (1-9)
    active_step: int = 1


class SessionCreateResponse(BaseModel):
    session_id: str


class ChatSendReq(BaseModel):
    session_id: str
    message: str
    active_step: Optional[int] = None


class ChatHistoryResp(BaseModel):
    session_id: str
    history: List[ChatTurn]


# ============================================================
# App + CORS
# ============================================================
app = FastAPI(title="Hopscotch IRML Chat API", version="0.2.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:9581", "http://127.0.0.1:9581",
        "https://hopscotchai.us", "https://hopscotch4all.com",
        # Native iPad/iOS app (Capacitor) origins
        "capacitor://localhost", "https://localhost", "ionic://localhost",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================
# Helpers: sessions
# ============================================================
def _require_session(session_id: str) -> SessionData:
    """Load a session from MongoDB and return it as a SessionData model."""
    doc = find_session(session_id)
    if not doc:
        raise HTTPException(status_code=404, detail="Session not found")
    chat_turns = [ChatTurn(**t) for t in (doc.get("chat") or [])]
    return SessionData(
        id=doc["session_id"],
        created_at=doc.get("created_at", ""),
        chat=chat_turns,
        worldview_band=doc.get("worldview_band"),
        worldview_label=doc.get("worldview_label"),
        step_notes=doc.get("step_notes", {}),
        resolved_path=doc.get("resolved_path"),
        chosen_methodology=doc.get("chosen_methodology"),
        active_step=doc.get("active_step", 1),
    )


def _persist_session(sess: SessionData):
    """Write the current SessionData back to MongoDB."""
    update_session(sess.id, {
        "chat": [t.dict() for t in sess.chat],
        "worldview_band": sess.worldview_band,
        "worldview_label": sess.worldview_label,
        "step_notes": sess.step_notes,
        "resolved_path": sess.resolved_path,
        "chosen_methodology": sess.chosen_methodology,
        "active_step": sess.active_step,
    })



WORLDVIEW_DESCRIPTIONS = {
    "positivist": (
        "Positivist: Believes in an objective, knowable reality. Knowledge is gained through "
        "observation, measurement, and empirical testing. Research should be value-free and "
        "generalizable. Most often associated with quantitative methods — experiments, surveys, "
        "statistical analysis. The researcher remains detached and neutral."
    ),
    "post_positivist": (
        "Post-Positivist: Acknowledges that reality exists but can only be imperfectly known. "
        "All observation is fallible and theory-laden. Emphasises falsification, triangulation, "
        "and critical multiplism. Most often associated with quantitative methods but recognises "
        "limitations of absolute objectivity. The researcher strives for objectivity while "
        "acknowledging bias."
    ),
    "constructivist": (
        "Constructivist (Interpretivist): Believes reality is socially constructed and that "
        "multiple, equally valid realities exist. Knowledge is co-created between researcher "
        "and participants. Values deep understanding of lived experiences, meaning-making, and "
        "context. Most often associated with qualitative methods — interviews, observations, "
        "narrative analysis — though constructivist researchers also legitimately conduct "
        "quantitative and mixed methods studies (e.g. validated surveys or quasi-experiments "
        "on constructivist learning environments). The researcher is an active participant "
        "in the research process."
    ),
    "transformative": (
        "Transformative: Centres issues of power, justice, and equity. Reality is shaped by "
        "social, political, cultural, and economic forces. Research should serve marginalised "
        "communities and promote social change. Most often associated with qualitative and "
        "participatory methods, though quantitative evidence is also used to expose inequities. "
        "The researcher is an advocate who collaborates with communities."
    ),
    "pragmatist": (
        "Pragmatist: Focuses on 'what works' rather than committing to a single ontology. "
        "The research question drives the choice of methods — quantitative, qualitative, or both. "
        "Values practical consequences, real-world applicability, and problem-solving. "
        "Embraces mixed methods and methodological flexibility. The researcher chooses approaches "
        "based on the nature of the problem being studied."
    ),
}


def _render_worldview_profile(sess: SessionData) -> str:
    """
    Human-readable summary sent to the LLM with rich worldview context.
    """
    if not sess.worldview_band:
        return "The student has not yet selected a worldview."
    band = sess.worldview_band
    label = sess.worldview_label or band.replace("_", " ").title()
    desc = WORLDVIEW_DESCRIPTIONS.get(band, "")
    path = sess.resolved_path or "not yet determined"
    chosen = sess.chosen_methodology
    parts = [f"Student's worldview: {label}"]
    if chosen and chosen != path:
        parts.append(
            f"Research methodology pathway: {chosen} (the student's own choice; "
            f"their worldview's usual pathway would be {path} - support their choice)"
        )
    else:
        parts.append(
            f"Research methodology pathway: {path} (the usual default for this worldview - "
            f"NOT a requirement; the student may choose a different pathway in Step 4)"
        )
    if desc:
        parts.append(f"Worldview description: {desc}")
    return "\n".join(parts)


def _get_chat(sess: SessionData) -> List[ChatTurn]:
    if sess.chat is None:
        sess.chat = []
    return sess.chat


# ============================================================
# RAG (with pdfminer + keyword fallback)
# ============================================================
def _read_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_md(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    txt = ""
    # Try pypdf first
    if PdfReader is not None:
        try:
            reader = PdfReader(str(path))
            txt = "\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception:
            txt = ""
    # Fallback to pdfminer if installed
    if (not txt or not txt.strip()) and pdfminer_extract_text is not None:
        try:
            txt = pdfminer_extract_text(str(path)) or ""
        except Exception:
            txt = ""
    return txt


def _load_all_docs() -> List[Dict[str, str]]:
    docs: List[Dict[str, str]] = []
    if not DOCS_DIR.exists():
        return docs
    for p in sorted(DOCS_DIR.glob("**/*")):
        ext = p.suffix.lower()
        try:
            if ext == ".txt":
                docs.append({"source": p.name, "text": _read_txt(p)})
            elif ext in (".md", ".markdown"):
                docs.append({"source": p.name, "text": _read_md(p)})
            elif ext == ".pdf":
                docs.append({"source": p.name, "text": _read_pdf(p)})
        except Exception:
            continue
    return docs


def _chunk(text: str, max_chars: int = 2400, overlap: int = 400) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    step = max(1, max_chars - overlap)
    return [text[i: i + max_chars] for i in range(0, len(text), step)]


def _ensure_embedder():
    global _embedder
    if not RAG_AVAILABLE:
        return
    if _embedder is None:
        _embedder = SentenceTransformer(EMBED_MODEL_NAME)


def _build_index(force: bool = False):
    """Build or load the FAISS index; if RAG unavailable, no-op.
    With force=True, always rebuild from the resources folder (used by the
    admin 'rebuild knowledge base' action after files change)."""
    global _faiss_index, _chunks, _raw_docs_cache
    if not RAG_AVAILABLE:
        _faiss_index = None
        _chunks = []
        return {"rag_available": False, "sources": 0, "chunks": 0}

    _ensure_embedder()

    if not force and INDEX_PATH.exists() and META_PATH.exists():
        try:
            _faiss_index = faiss.read_index(str(INDEX_PATH))
            _chunks = json.loads(META_PATH.read_text(encoding="utf-8"))
            return {"rag_available": True,
                    "sources": len({c.get("source") for c in _chunks}),
                    "chunks": len(_chunks)}
        except Exception as e:
            logger.warning("Failed to load existing index; rebuilding. %s", e)

    # Fresh build — re-read every doc so newly added/removed files are reflected.
    _raw_docs_cache = None
    docs = _load_all_docs()
    chunks: List[Dict[str, Any]] = []
    for d in docs:
        for piece in _chunk(d["text"]):
            chunks.append({"text": piece, "source": d["source"]})

    if not chunks:
        _faiss_index = faiss.IndexFlatIP(EMBED_DIM) if RAG_AVAILABLE else None
        _chunks = []
        return {"rag_available": True, "sources": 0, "chunks": 0}

    texts = [c["text"] for c in chunks]
    _ensure_embedder()
    vecs = _embedder.encode(
        texts, convert_to_numpy=True, normalize_embeddings=True
    )

    index = faiss.IndexFlatIP(vecs.shape[1])
    index.add(vecs)
    _faiss_index = index

    faiss.write_index(_faiss_index, str(INDEX_PATH))
    META_PATH.write_text(
        json.dumps(
            [{"id": i, "text": c["text"], "source": c["source"]} for i, c in enumerate(chunks)],
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    _chunks = [{"id": i, **c} for i, c in enumerate(chunks)]
    return {"rag_available": True,
            "sources": len({c["source"] for c in chunks}),
            "chunks": len(chunks)}


# Extensions the knowledge base can ingest.
RESOURCE_EXTS = (".pdf", ".txt", ".md", ".markdown")


def _list_resource_files() -> List[Dict[str, Any]]:
    """List files in the resources folder with index coverage info."""
    indexed_counts: Dict[str, int] = {}
    for c in _chunks:
        s = c.get("source", "")
        indexed_counts[s] = indexed_counts.get(s, 0) + 1
    files: List[Dict[str, Any]] = []
    if DOCS_DIR.exists():
        for p in sorted(DOCS_DIR.glob("*")):
            if not p.is_file() or p.suffix.lower() not in RESOURCE_EXTS:
                continue
            st = p.stat()
            files.append({
                "name": p.name,
                "ext": p.suffix.lower().lstrip("."),
                "size": st.st_size,
                "modified": datetime.utcfromtimestamp(st.st_mtime).isoformat() + "Z",
                "chunks": indexed_counts.get(p.name, 0),
                "indexed": p.name in indexed_counts,
            })
    return files


def _index_stale() -> bool:
    """True if the resources folder no longer matches the built index."""
    on_disk = {p.name for p in DOCS_DIR.glob("*")
               if p.is_file() and p.suffix.lower() in RESOURCE_EXTS} if DOCS_DIR.exists() else set()
    in_index = {c.get("source", "") for c in _chunks}
    return on_disk != in_index


def _keyword_fallback(query: str, k: int = 5) -> List[Dict[str, Any]]:
    """Very simple keyword scoring fallback when FAISS/chunks unavailable."""
    global _raw_docs_cache
    if _raw_docs_cache is None:
        _raw_docs_cache = _load_all_docs()
    q = (query or "").strip().lower()
    if not q:
        return []
    scored: List[Dict[str, Any]] = []
    for d in _raw_docs_cache:
        text = d.get("text") or ""
        if not text:
            continue
        tl = text.lower()
        occ = tl.count(q)
        score = occ + (1.0 if q in tl else 0.0)
        if score > 0:
            scored.append(
                {
                    "text": text[:2000],
                    "source": d["source"],
                    "score": float(score),
                }
            )
    scored.sort(key=lambda x: x["score"], reverse=True)
    return scored[:k]


def _retrieve(query: str, k: int = 5) -> List[Dict[str, Any]]:
    """Try vector search; if nothing, use keyword fallback."""
    if RAG_AVAILABLE and _faiss_index is not None and _chunks:
        _ensure_embedder()
        try:
            qv = _embedder.encode(
                [query], convert_to_numpy=True, normalize_embeddings=True
            )
            D, I = _faiss_index.search(qv, k)
            out = []
            for idx, score in zip(I[0], D[0]):
                idx = int(idx)
                if 0 <= idx < len(_chunks):
                    ch = _chunks[idx]
                    out.append(
                        {
                            "text": ch["text"],
                            "source": ch["source"],
                            "score": float(score),
                        }
                    )
            if out:
                return out
        except Exception as e:
            logger.warning(
                "Vector retrieval failed; falling back to keywords. %s", e
            )
    return _keyword_fallback(query, k=k)


# ---- GeoIP resolution (ip-api.com, free, no key) ----

_GEO_CACHE: Dict[str, Dict] = {}

def _resolve_geo(ip: str) -> dict:
    """Resolve IP to geo data. Returns {} on failure."""
    if not ip or ip in ("127.0.0.1", "::1", "localhost"):
        return {}
    if ip in _GEO_CACHE:
        return _GEO_CACHE[ip]
    try:
        resp = requests.get(
            f"http://ip-api.com/json/{ip}?fields=status,city,regionName,country,lat,lon",
            timeout=3,
        )
        data = resp.json()
        if data.get("status") == "success":
            geo = {
                "city": data.get("city", ""),
                "regionName": data.get("regionName", ""),
                "country": data.get("country", ""),
                "lat": data.get("lat"),
                "lng": data.get("lon"),
            }
            _GEO_CACHE[ip] = geo
            return geo
    except Exception as e:
        logger.warning("GeoIP lookup failed for %s: %s", ip, e)
    return {}


def _extract_client_ip(request: Request) -> str:
    """Extract real client IP, accounting for Cloudflare / proxies."""
    return (
        request.headers.get("CF-Connecting-IP")
        or request.headers.get("X-Forwarded-For", "").split(",")[0].strip()
        or (request.client.host if request.client else "")
    )


# ---- Admin seed from env vars ----

ADMIN_EMAIL = os.environ.get("ADMIN_EMAIL")
ADMIN_PASSWORD = os.environ.get("ADMIN_PASSWORD")

def _seed_admin():
    """Auto-create admin user from env vars on startup if not exists."""
    if not ADMIN_EMAIL or not ADMIN_PASSWORD:
        logger.info("ADMIN_EMAIL/ADMIN_PASSWORD not set — skipping admin seed.")
        return
    existing = find_user_by_email(ADMIN_EMAIL)
    if existing:
        if existing.get("role") != "admin":
            from database import update_user_fields
            update_user_fields(str(existing["_id"]), {"role": "admin"})
            logger.info("Promoted existing user %s to admin.", ADMIN_EMAIL)
        else:
            logger.info("Admin user %s already exists.", ADMIN_EMAIL)
        return
    pw_hash = hash_password(ADMIN_PASSWORD)
    create_user(ADMIN_EMAIL, pw_hash, "admin", "Administrator", "higher_ed")
    logger.info("Created admin user: %s", ADMIN_EMAIL)


def _seed_glossary():
    """Populate the glossary collection from glossary_seed.json on first run."""
    try:
        seed_path = ROOT / "glossary_seed.json"
        if not seed_path.exists():
            return
        with open(seed_path, "r", encoding="utf-8") as f:
            terms = json.load(f)
        inserted = seed_glossary_if_empty(terms)
        if inserted:
            print(f"[glossary] Seeded {inserted} terms into the database")
    except Exception as e:
        print(f"[glossary] Seed skipped: {e}")


def _seed_step_resources():
    """Populate step resources (video + interactive per step/level) on first run."""
    try:
        seed_path = ROOT / "step_resources_seed.json"
        if not seed_path.exists():
            return
        with open(seed_path, "r", encoding="utf-8") as f:
            seed = json.load(f)
        inserted = seed_step_resources_if_empty(seed)
        if inserted:
            print(f"[step-resources] Seeded {inserted} rows into the database")
    except Exception as e:
        print(f"[step-resources] Seed skipped: {e}")


@app.on_event("startup")
def _startup():
    ensure_indexes()
    _build_index()
    load_paths_config()
    _seed_admin()
    _seed_glossary()
    _seed_step_resources()
    # Pre-warm the LLM so the first chat request doesn't cold-start
    _warm_llm()


def _warm_llm():
    """Pre-warm the LLM — works with both vLLM and Ollama backends."""
    if LLM_BACKEND == "vllm":
        try:
            logger.info("Pre-warming vLLM model %s ...", VLLM_MODEL)
            headers = {"Content-Type": "application/json"}
            if VLLM_API_KEY:
                headers["Authorization"] = f"Bearer {VLLM_API_KEY}"
            resp = requests.post(VLLM_URL, json={
                "model": VLLM_MODEL,
                "messages": [{"role": "user", "content": "hi"}],
                "max_tokens": 1,
            }, headers=headers, timeout=300)
            resp.raise_for_status()
            logger.info("vLLM model %s is warm and ready.", VLLM_MODEL)
        except Exception as e:
            logger.warning("Failed to pre-warm vLLM model: %s — will try Ollama fallback", e)
    else:
        try:
            logger.info("Pre-warming Ollama model %s ...", LLM_MODEL)
            resp = requests.post(OLLAMA_URL, json={
                "model": LLM_MODEL,
                "messages": [{"role": "user", "content": "hi"}],
                "stream": False,
                "options": {"num_predict": 1},
            }, timeout=300)
            resp.raise_for_status()
            logger.info("Ollama model %s is warm and ready.", LLM_MODEL)
        except Exception as e:
            logger.warning("Failed to pre-warm Ollama model: %s", e)


# ============================================================
# LLM call (Ollama)
# ============================================================

def _get_step_llm_guidance(sess: SessionData, active_step: Optional[int]) -> Optional[str]:
    """Resolve the LLM guidance string for the given step from the paths config."""
    if not active_step or active_step < 4:
        return None
    paths_cfg = load_paths_config()
    resolved = sess.resolved_path
    if not resolved:
        return None

    all_paths = paths_cfg.get("paths", {})
    path_data = all_paths.get(resolved, {})
    step_cfg = path_data.get("steps", {}).get(str(active_step), {})

    # Handle mixed-methods inheritance for steps 5-9
    if resolved == "mixed" and active_step >= 5 and step_cfg.get("inherits_from_chosen_methodology"):
        chosen = sess.chosen_methodology
        if not chosen:
            return "The student has not yet chosen their primary methodology in Step 4."
        inherited = all_paths.get(chosen, {}).get("steps", {}).get(str(active_step), {})
        guidance = inherited.get("llm_guidance", "")
        addendum = step_cfg.get("llm_guidance_addendum", "")
        if addendum:
            guidance = f"{guidance}\n{addendum}" if guidance else addendum
        return guidance or None

    # Non-mixed path but student overrode methodology at Step 4
    if resolved != "mixed" and sess.chosen_methodology and sess.chosen_methodology != resolved:
        override_cfg = all_paths.get(sess.chosen_methodology, {}).get("steps", {}).get(str(active_step), {})
        if override_cfg:
            return override_cfg.get("llm_guidance")

    return step_cfg.get("llm_guidance")


STEP_NAMES = {
    1: "Worldview / Paradigm",
    2: "Topic & Goals",
    3: "Literature & Conceptual Framework",
    4: "Methodology / Research Design",
    5: "Research Question(s)",
    6: "Data Collection",
    7: "Data Analysis",
    8: "Trustworthiness / Validity",
    9: "Ethics",
}


def build_ollama_payload(worldview_profile, step_context, user_msg, passages,
                         stream=False, active_step=None, step_llm_guidance=None,
                         chat_history=None, language="en"):
    """
    Shared helper to build the Ollama chat payload.
    Set stream=True when you want chunked responses, False for normal JSON.
    """
    ctx_blocks = []
    for i, p in enumerate(passages):
        ctx_blocks.append(
            f"[{i+1}] Source: {p['source']}\n{p['text'][:800]}"
        )
    ctx_text = "\n\n".join(ctx_blocks) if ctx_blocks else "No matching passages."

    system_msg = (
        "You are a knowledgeable, supportive research-methods tutor embedded in the "
        "Hopscotch IRML (Introductory Research Methods Learning) platform. You help "
        "students scaffold their research design through a 9-step process.\n\n"

        "==================================================================\n"
        "CORE RULE — YOU ARE A COACH, NOT A GHOSTWRITER (applies to ALL 9 steps)\n"
        "==================================================================\n"
        "Your job is to help students LEARN to design research — NOT to produce their "
        "research design for them. This is your single most important rule, and it "
        "overrides any request to the contrary:\n"
        "- ALLOWED: explain concepts and terminology, illustrate an idea with a "
        "general example, answer 'what does X mean?' questions, ask guiding questions, "
        "and give specific feedback on content the student has ALREADY written.\n"
        "- NOT ALLOWED: writing or generating the student's OWN design content for "
        "them — their research topic, research question, aim, hypothesis, literature or "
        "theoretical framework, methodology choice, data-collection plan, analysis plan, "
        "trustworthiness plan, or ethics plan. Never hand over ready-to-paste answers, "
        "even if the student asks directly, repeatedly, or frames it as 'just an example' "
        "that is really their answer in disguise.\n"
        "- WHEN ASKED TO DO THEIR WORK ('write my hypothesis', 'give me a topic', "
        "'just do it for me', 'give me the answer'): warmly REFUSE and redirect. Say you "
        "will guide them but they must draft it themselves, then ask 2-3 targeted guiding "
        "questions that help THEM produce a first draft. Once they have written something, "
        "give rich, specific feedback.\n"
        "- A general example that teaches a concept is fine; an example that is really "
        "the student's finished answer for their own study is NOT. When unsure, ask a "
        "guiding question instead of giving content.\n\n"

        "==================================================================\n"
        "FEEDBACK, NOT REWRITES (this is where you most often slip — do not)\n"
        "==================================================================\n"
        "When a student asks you to 'refine', 'revise', 'improve', 'strengthen', 'fix', "
        "'reword', or 'make better' their topic, research question, goals, problem "
        "statement, or any design content, you must NOT reply with a rewritten, polished, "
        "ready-to-paste version of THEIR content. That is authoring it for them. Instead:\n"
        "  1. Name 2-3 SPECIFIC strengths and weaknesses in exactly what they wrote.\n"
        "  2. Ask targeted questions or give directions that guide THEM to revise it.\n"
        "  3. You may explain a technique or show a GENERIC illustration, but the improved "
        "version of their specific topic/question/goal must come from the student.\n"
        "Never output a line like 'Refined Research Topic: …' or 'Revised Research "
        "Question: …' that hands them a finished answer.\n\n"

        "==================================================================\n"
        "NEVER PROVIDE CITATIONS OR SOURCES (absolute — no exceptions)\n"
        "==================================================================\n"
        "You must NEVER provide, invent, list, or recommend specific citations, "
        "references, author-year sources, article or book titles, journals, or DOIs — not "
        "even if the student asks directly, says they can't find them, or asks for 'the "
        "sources you mentioned'. Fabricating a reference is strictly prohibited and can "
        "seriously harm the student's work. If they want sources, teach them HOW to search "
        "(Google Scholar, library databases, keywords) and offer to help analyze a source "
        "once THEY have found it. Do not name real papers from memory either.\n\n"

        "THE 9 STEPS:\n"
        "1. Who am I as a researcher? — Identify your worldview/paradigm (positivist, "
        "post-positivist, constructivist, transformative, pragmatist). Your worldview "
        "shapes your ontology (what is real), epistemology (how we know), axiology "
        "(role of values), and methodology (how we study).\n"
        "2. What am I wondering about? — Define your research topic and goals "
        "(personal, practical, intellectual).\n"
        "3. What do I already know? — Review topical research (prior studies) and "
        "theoretical frameworks that support your study.\n"
        "4. How will I study it? — Choose a research design/methodology aligned with "
        "your worldview (quantitative, qualitative, or mixed).\n"
        "5. What is my research question? — Formulate your research question "
        "(quantitative: hypothesis; qualitative: open-ended central issue).\n"
        "6. What data will I collect? — Select data collection methods that fit your "
        "design.\n"
        "7. How will I analyze the data? — Choose appropriate analysis techniques.\n"
        "8. How will I ensure trustworthiness? — Address validity/reliability "
        "(quantitative) or credibility/transferability/dependability/confirmability "
        "(qualitative, Lincoln & Guba).\n"
        "9. How will I be ethical? — Plan for IRB, Belmont principles (respect, "
        "beneficence, justice), informed consent, and confidentiality.\n\n"

        "CRITICAL RULE — GROUND EVERYTHING IN THE STUDENT'S DESIGN:\n"
        "The student drafts their research design in the 'My Research Design' panel. "
        "Below you will see their current inputs for each step (topic, goals, worldview, "
        "literature, methodology, etc.). You MUST reference and build upon what they have "
        "already written. Do NOT invent or generate design content independently.\n"
        "- If the student has filled in a field, refer to their SPECIFIC inputs by name "
        "(e.g. 'Your topic about X…', 'Since you chose the pragmatist worldview…') and "
        "help them refine, strengthen, or expand what they wrote.\n"
        "- If a field is empty, encourage the student to write their initial thoughts in "
        "the 'My Research Design' panel first, then come back for feedback.\n"
        "- Never produce a full research design from scratch — your role is to coach and "
        "give feedback on what the student has drafted, not to do the work for them.\n\n"

        "QUESTION-DRIVEN COACHING (Steps 2, 3, and 4):\n"
        "For Steps 2, 3, and 4 you MUST be question-driven. Do NOT suggest or recommend "
        "specific research topics, literature, theoretical frameworks, or methodologies. "
        "The student must come up with their own ideas first.\n"
        "- Step 2 (Topic & Goals): Do NOT suggest topics. Ask guiding questions like "
        "'What issues in your field interest you the most?', 'What problem have you "
        "observed that you want to explore?', 'What would you like to change or understand "
        "better?' — let the student discover their own topic through reflection.\n"
        "- Step 3 (Literature Review): Do NOT recommend specific studies, authors, or "
        "frameworks. Instead ask 'What research have you already read on this topic?', "
        "'What theories from your coursework connect to your topic?', 'What gaps have you "
        "noticed in the existing research?' — guide them to identify their own sources.\n"
        "- Step 4 (Methodology): Do NOT prescribe a methodology. Ask 'Based on your "
        "worldview, what approach feels most natural?', 'Are you trying to measure "
        "something or understand experiences?', 'What type of data would best answer your "
        "question?' — let the student reason toward a methodology.\n"
        "- Once the student HAS written something in their design, THEN you may give "
        "substantive feedback, point out strengths, identify gaps, and suggest refinements. "
        "But always wait for their input first.\n\n"

        "YOUR APPROACH:\n"
        "- When explaining a worldview (Step 1), be substantive: discuss its ontology, "
        "epistemology, axiology, and methodology implications with concrete examples.\n"
        "- Worldviews are lenses, NOT locks: the research question ultimately drives the "
        "choice of methodology. A constructivist may legitimately conduct quantitative "
        "research (e.g. validated surveys, quasi-experiments on constructivist learning "
        "environments) and a post-positivist may conduct qualitative research. When a "
        "student asks whether they can use a methodology that differs from their "
        "worldview's usual pairing, or resists the usual pairing, do NOT insist on the "
        "default: acknowledge the legitimacy of their direction, present the options with "
        "their trade-offs, and let the student reason to their own choice. Tell them they "
        "can switch their methodology pathway in Step 4 if they wish.\n"
        "- For Steps 2-4: lead with guiding questions, then give feedback ONLY after "
        "the student has written their own content in 'My Research Design'.\n"
        "- For Steps 5-9: be substantive in EXPLAINING concepts and giving feedback on "
        "what the student wrote — but per the CORE RULE, still never author their research "
        "question, hypothesis, data-collection, analysis, trustworthiness, or ethics "
        "content for them. Guide them to write it; then critique and refine it.\n"
        "- Reference specific methodologies, frameworks, and scholars when relevant "
        "(only in response to what the student has already written, not as suggestions).\n"
        "- Use a warm, encouraging tone — the student may be new to research.\n"
        "- This is a chat, so a letter-style sign-off usually isn't needed. If you do "
        "close with one, sign ONLY as 'Hopscotch' — never a placeholder like '[Your "
        "Tutor]' or '[Your Name]'.\n"
        "- Keep responses focused but thorough (2-4 paragraphs typically).\n"
        "- Do NOT include article citations, source lists, or reference sections in your responses.\n"
        "- For Steps 1, 2, and 3: the student must find their own sources — do not suggest any.\n"
    )

    # Inject step-specific guidance + a hard "current-step lock" so the model
    # coaches only on the current step and never works ahead into later steps.
    if active_step:
        cur = int(active_step)
        cur_name = STEP_NAMES.get(cur, "")
        future = ", ".join(f"Step {n} ({STEP_NAMES[n]})" for n in range(cur + 1, 10))
        system_msg += (
            "\n==================================================================\n"
            f"CURRENT-STEP LOCK — the student is on STEP {cur}: {cur_name}\n"
            "==================================================================\n"
            f"Coach ONLY on Step {cur}. You may use everything the student has already "
            f"written in Steps 1-{cur} as context to give feedback on Step {cur}.\n"
            f"Steps {cur+1}-9 have NOT been reached yet and are OFF LIMITS:\n"
            f"- Do NOT write, draft, suggest, outline, give examples of, or work out ANY "
            f"content for later steps ({future or 'none — this is the last step'}). No "
            f"research questions, methodology, data-collection plans, analysis plans, "
            f"trustworthiness or ethics content — none of it — until the student is actually "
            f"on that step.\n"
            f"- Do NOT preview 'what comes next' with specifics. At most, a single line like "
            f"'you'll work that out in Step X', then bring them back to Step {cur}.\n"
            f"- If the student asks you to help with or answer a later step, warmly decline and "
            f"refocus them on Step {cur} — the point is for them to think each step through "
            f"themselves, in order.\n"
        )
        if step_llm_guidance:
            system_msg += f"\nStep-specific instructions for Step {cur}:\n{step_llm_guidance}\n"

    # Steps 1-3: no resource snippets — student must find their own sources
    if active_step and active_step <= 3:
        context_msg = (
            f"Student context:\n{worldview_profile}\n\n"
            f"STUDENT'S RESEARCH DESIGN (from 'My Research Design' panel — reference these directly):\n{step_context}\n"
        )
    else:
        context_msg = (
            f"Student context:\n{worldview_profile}\n\n"
            f"STUDENT'S RESEARCH DESIGN (from 'My Research Design' panel — reference these directly):\n{step_context}\n\n"
            f"IRML resource snippets:\n{ctx_text}"
        )

    # Build messages: system + context + conversation history + latest user msg
    if language == "es":
        system_msg += (
            "\n\nLANGUAGE: The student uses Hopscotch in Spanish. ALWAYS respond "
            "entirely in Spanish, with a warm academic tone (use 'cosmovisión' for "
            "worldview). Keep methodology terminology accurate in Spanish."
        )
    elif language == "zh":
        system_msg += (
            "\n\nLANGUAGE: The student uses Hopscotch in Chinese. ALWAYS respond "
            "entirely in Simplified Chinese (简体中文), with a warm academic tone. "
            "Use standard research-methods terminology: 定量研究 for quantitative "
            "research, 定性研究 for qualitative research, 混合研究方法 for mixed "
            "methods, 世界观 for worldview."
        )

    messages = [
        {"role": "system", "content": system_msg},
        {"role": "system", "content": context_msg},
    ]

    # Include recent conversation history (last 20 turns to stay within context)
    if chat_history:
        recent = chat_history[-20:]
        for turn in recent:
            # Skip the very last user message — we append it separately below
            if turn is recent[-1] and turn.role == "user" and turn.content == user_msg:
                continue
            content = turn.content
            # Strip any legacy "Quick references" sections from older assistant messages
            if turn.role == "assistant" and content:
                marker = "**Quick references from our notes:**"
                ix = content.find(marker)
                if ix != -1:
                    content = content[:ix].rstrip()
            messages.append({"role": turn.role, "content": content})

    messages.append({"role": "user", "content": user_msg})

    return {
        "model": LLM_MODEL,
        "stream": stream,
        "options": {"temperature": LLM_TEMP},
        "messages": messages,
    }


def _call_vllm(messages: list, temperature: float = LLM_TEMP,
               max_tokens: int = 2048, timeout: int = 120) -> Optional[str]:
    """Call vLLM (OpenAI-compatible API). Returns content string or None on failure."""
    headers = {"Content-Type": "application/json"}
    if VLLM_API_KEY:
        headers["Authorization"] = f"Bearer {VLLM_API_KEY}"
    try:
        resp = requests.post(VLLM_URL, json={
            "model": VLLM_MODEL,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_tokens,
            "stream": False,
        }, headers=headers, timeout=timeout)
        resp.raise_for_status()
        data = resp.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        logger.warning("vLLM call failed: %s", e)
        return None


def _call_ollama(payload: dict, timeout: int = 120) -> Optional[str]:
    """Call Ollama. Returns content string or None on failure."""
    try:
        resp = requests.post(OLLAMA_URL, json=payload, timeout=timeout)
        resp.raise_for_status()
        data = resp.json()
        return data.get("message", {}).get("content", "").strip()
    except Exception as e:
        logger.warning("Ollama call failed: %s", e)
        return None


def call_llm(worldview_profile: str, step_context: str, user_msg: str,
             passages: List[Dict[str, Any]],
             active_step: Optional[int] = None,
             step_llm_guidance: Optional[str] = None,
             chat_history=None, language: str = "en") -> str:

    payload = build_ollama_payload(
        worldview_profile, step_context, user_msg, passages,
        stream=False, active_step=active_step, step_llm_guidance=step_llm_guidance,
        chat_history=chat_history, language=language,
    )

    result = None
    # Try vLLM first if configured
    if LLM_BACKEND == "vllm":
        result = _call_vllm(payload["messages"], temperature=LLM_TEMP)
        if result is None:
            logger.info("vLLM failed, falling back to Ollama...")
            result = _call_ollama(payload)
    else:
        result = _call_ollama(payload)
        if result is None and LLM_BACKEND != "vllm":
            logger.info("Ollama failed, trying vLLM fallback...")
            result = _call_vllm(payload["messages"], temperature=LLM_TEMP)

    return result or (
        "I ran into an issue calling the language model. "
        "Please try again or check the backend logs."
    )


# ============================================================
# Auth endpoints
# ============================================================

class RegisterReq(BaseModel):
    email: str
    password: str
    name: str
    role: str  # "student" or "teacher"
    education_level: str = "high_school"  # "high_school" or "higher_ed"


class LoginReq(BaseModel):
    email: str
    password: str


class AuthResp(BaseModel):
    token: str
    email: Optional[str] = None
    username: Optional[str] = None
    name: str
    role: str
    education_level: str = "high_school"
    ai_enabled: bool = True
    access_mode: str = "full"
    unlocked_phase: Optional[int] = None


def _student_class_settings(user: dict) -> dict:
    """Effective teacher-controlled settings for this user's class. Non-classroom
    users (faculty, email students, teachers, admins) get the permissive defaults
    (AI on, full access)."""
    if not user or user.get("role") != "classroom_student":
        return get_class_settings(None)
    class_id = user.get("class_id")
    if not class_id:
        return get_class_settings(None)
    return get_class_settings(find_class_by_id(str(class_id)))


def _student_ai_enabled(user: dict) -> bool:
    """Whether the AI assistant is available to this user (teacher-controlled)."""
    return bool(_student_class_settings(user).get("ai_enabled", True))


# Access/pacing modes (Phase 2). The 9 steps are grouped into 3 phases.
STEP_PHASES = [[1, 2, 3], [4, 5, 6], [7, 8, 9]]


def _phase_of_step(step: int) -> int:
    for i, ph in enumerate(STEP_PHASES):
        if step in ph:
            return i + 1
    return 1


def _step_is_locked(step: int, settings: dict, completed_steps) -> bool:
    """Whether a step is locked for a student, per the class access mode.
    - full: nothing locked
    - step: sequential — you may work on completed steps and the first incomplete
      one, but not jump ahead
    - phase: steps in phases beyond the teacher-unlocked phase are locked"""
    mode = (settings or {}).get("access_mode", "full")
    completed = set(completed_steps or [])
    if mode == "step":
        if step <= 1 or step in completed:
            return False
        first_incomplete = next((n for n in range(1, 10) if n not in completed), 10)
        return step > first_incomplete
    if mode == "phase":
        unlocked = (settings or {}).get("unlocked_phase") or 1
        return _phase_of_step(step) > unlocked
    return False


def _guard_step_access(user: dict, sess: "SessionData", step: int):
    """Raise 403 if the student's class access mode locks this step."""
    settings = _student_class_settings(user)
    if settings.get("access_mode", "full") == "full":
        return
    completed = _compute_completed_steps_from_session(sess)
    if _step_is_locked(step, settings, completed):
        raise HTTPException(status_code=403, detail="This step is locked by your teacher.")


@app.post("/auth/register", response_model=AuthResp)
def register(req: RegisterReq):
    if req.role not in ("student", "teacher"):
        raise HTTPException(status_code=400, detail="Role must be 'student' or 'teacher'")
    if req.education_level not in ("high_school", "higher_ed"):
        raise HTTPException(status_code=400, detail="education_level must be 'high_school' or 'higher_ed'")
    if find_user_by_email(req.email):
        raise HTTPException(status_code=409, detail="Email already registered")
    pw_hash = hash_password(req.password)
    create_user(req.email, pw_hash, req.role, req.name, req.education_level)
    token = create_access_token({"sub": req.email})
    return AuthResp(token=token, email=req.email, name=req.name,
                    role=req.role, education_level=req.education_level)


@app.post("/auth/login", response_model=AuthResp)
def login(req: LoginReq, request: Request):
    client_ip = _extract_client_ip(request)
    user_agent = request.headers.get("User-Agent", "")
    user = find_user_by_email(req.email)
    if not user or not verify_password(req.password, user["password_hash"]):
        if user:
            record_login(str(user["_id"]), req.email, client_ip, {}, user_agent, success=False)
        raise HTTPException(status_code=401, detail="Invalid email or password")
    if not user.get("is_active", True):
        raise HTTPException(status_code=403, detail="Account deactivated. Contact your administrator.")
    geo = _resolve_geo(client_ip)
    record_login(str(user["_id"]), req.email, client_ip, geo, user_agent, success=True)
    update_user_fields(str(user["_id"]), {
        "last_login_at": datetime.utcnow().isoformat() + "Z",
        "last_login_ip": client_ip,
    })
    token = create_access_token({"sub": req.email})
    return AuthResp(token=token, email=req.email, name=user["name"],
                    role=user["role"],
                    education_level=user.get("education_level", "high_school"))


class LanguageReq(BaseModel):
    language: str


SUPPORTED_LANGUAGES = {"en", "es", "zh"}


@app.post("/auth/language")
def set_language(req: LanguageReq, user: dict = Depends(get_current_user)):
    """Persist the user's interface language (Phase 1: en / es)."""
    lang = (req.language or "").strip().lower()
    if lang not in SUPPORTED_LANGUAGES:
        raise HTTPException(status_code=400, detail="Unsupported language")
    update_user_fields(str(user["_id"]), {"language": lang})
    return {"ok": True, "language": lang}


@app.get("/auth/me")
def get_me(user: dict = Depends(get_current_user)):
    return {
        "email": user.get("email"),
        "username": user.get("username"),
        "name": user["name"],
        "role": user["role"],
        "education_level": user.get("education_level", "high_school"),
        "language": user.get("language", "en"),
        "ai_enabled": _student_ai_enabled(user),
        "access_mode": _student_class_settings(user).get("access_mode", "full"),
        "unlocked_phase": _student_class_settings(user).get("unlocked_phase"),
    }


@app.post("/auth/refresh")
def refresh_token(user: dict = Depends(get_current_user)):
    """Sliding session: while the current token is still valid, hand back a
    fresh 7-day token plus the current profile (same shape as /auth/me). The
    client swaps it in on load / tab refocus, so active users are never forced
    to log back in. An expired token fails get_current_user with 401, which the
    client treats as a real logout."""
    if user.get("email"):
        token = create_access_token({"sub": user["email"]})
    else:
        token = create_access_token({"sub": user["username"], "sub_type": "username"})
    return {
        "token": token,
        "email": user.get("email"),
        "username": user.get("username"),
        "name": user["name"],
        "role": user["role"],
        "education_level": user.get("education_level", "high_school"),
        "language": user.get("language", "en"),
        "ai_enabled": _student_ai_enabled(user),
        "access_mode": _student_class_settings(user).get("access_mode", "full"),
        "unlocked_phase": _student_class_settings(user).get("unlocked_phase"),
    }


# ---------- Password reset ----------

class ForgotPasswordReq(BaseModel):
    email: str

class ResetPasswordReq(BaseModel):
    token: str
    new_password: str


def _send_reset_email(to_email: str, reset_token: str):
    """Send a password-reset link via Resend. Fails silently to avoid leaking user existence."""
    if not RESEND_API_KEY:
        logger.warning("RESEND_API_KEY not configured — cannot send reset email")
        return

    resend.api_key = RESEND_API_KEY
    reset_link = f"{FRONTEND_URL}?reset_token={reset_token}"

    html_body = f"""\
<div style="font-family: sans-serif; max-width: 480px; margin: 0 auto;">
    <h2 style="color: #2B5EA7;">Hopscotch Password Reset</h2>
    <p>You requested a password reset for your Hopscotch account.</p>
    <p>
        <a href="{reset_link}"
           style="display: inline-block; padding: 12px 24px; background: #2B5EA7;
                  color: white; text-decoration: none; border-radius: 8px; font-weight: 600;">
            Reset Password
        </a>
    </p>
    <p style="font-size: 0.85rem; color: #666;">
        This link is valid for 30 minutes. If you did not request this, ignore this email.
    </p>
</div>"""

    try:
        resend.Emails.send({
            "from": f"Hopscotch <{EMAIL_FROM}>",
            "to": [to_email],
            "subject": "Hopscotch - Reset Your Password",
            "html": html_body,
        })
        logger.info(f"Reset email sent to {to_email}")
    except Exception as e:
        logger.error(f"Failed to send reset email: {e}")


@app.post("/auth/forgot-password")
def forgot_password(req: ForgotPasswordReq):
    """Send a password-reset email. Always returns success to avoid leaking user existence."""
    user = find_user_by_email(req.email)
    if user:
        token = create_password_reset_token(req.email)
        _send_reset_email(req.email, token)
    return {"message": "If an account with that email exists, a reset link has been sent."}


@app.post("/auth/reset-password")
def reset_password(req: ResetPasswordReq):
    """Verify the reset token and update the user's password."""
    payload = decode_token(req.token)

    if payload.get("purpose") != "password_reset":
        raise HTTPException(status_code=400, detail="Invalid reset token")

    email = payload.get("sub")
    if not email:
        raise HTTPException(status_code=400, detail="Invalid reset token")

    user = find_user_by_email(email)
    if not user:
        raise HTTPException(status_code=400, detail="Invalid reset token")

    if len(req.new_password) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters")

    new_hash = hash_password(req.new_password)
    update_user_password(email, new_hash)

    return {"message": "Password updated successfully. You can now log in."}


# ---------- Classroom login ----------

class ClassroomLoginReq(BaseModel):
    username: str
    password: str


@app.post("/auth/classroom-login", response_model=AuthResp)
def classroom_login(req: ClassroomLoginReq, request: Request):
    """Login for classroom (username-based) students."""
    client_ip = _extract_client_ip(request)
    user_agent = request.headers.get("User-Agent", "")
    user = find_user_by_username(req.username)
    if not user or user.get("role") != "classroom_student":
        raise HTTPException(status_code=401, detail="Invalid username or password")
    if not verify_password(req.password, user["password_hash"]):
        record_login(str(user["_id"]), req.username, client_ip, {}, user_agent, success=False)
        raise HTTPException(status_code=401, detail="Invalid username or password")
    if not user.get("is_active", True):
        raise HTTPException(status_code=403, detail="Account deactivated. Contact your administrator.")
    geo = _resolve_geo(client_ip)
    record_login(str(user["_id"]), req.username, client_ip, geo, user_agent, success=True)
    update_user_fields(str(user["_id"]), {
        "last_login_at": datetime.utcnow().isoformat() + "Z",
        "last_login_ip": client_ip,
    })
    token = create_access_token({"sub": req.username, "sub_type": "username"})
    _cs = _student_class_settings(user)
    return AuthResp(
        token=token,
        username=req.username,
        name=user["name"],
        role=user["role"],
        education_level=user.get("education_level", "high_school"),
        ai_enabled=bool(_cs.get("ai_enabled", True)),
        access_mode=_cs.get("access_mode", "full"),
        unlocked_phase=_cs.get("unlocked_phase"),
    )


# ============================================================
# Teacher endpoints
# ============================================================

class CreateClassReq(BaseModel):
    class_name: str
    student_count: int
    password: str


@app.post("/teacher/create-class")
def create_class_endpoint(req: CreateClassReq, user: dict = Depends(get_current_user)):
    if user.get("role") != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can create classes")
    if not (1 <= req.student_count <= 100):
        raise HTTPException(status_code=400, detail="Student count must be between 1 and 100")
    if len(req.password) < 4:
        raise HTTPException(status_code=400, detail="Password must be at least 4 characters")

    # Generate class_code from class_name
    raw = re.sub(r'[^a-z0-9]', '', req.class_name.lower().replace(' ', ''))
    class_code = raw[:20] or "class"

    # Ensure uniqueness
    base_code = class_code
    counter = 1
    while find_class_by_code(class_code):
        class_code = f"{base_code}{counter}"
        counter += 1

    pw_hash = hash_password(req.password)
    teacher_id = str(user["_id"])
    teacher_edu_level = user.get("education_level", "high_school")

    class_id = create_class_doc(teacher_id, req.class_name, class_code, pw_hash, req.password, req.student_count)

    # Create student accounts — inherit teacher's education_level
    students = []
    for i in range(1, req.student_count + 1):
        username = f"{class_code}_{i:02d}"
        student_name = f"Student {i:02d}"
        create_classroom_student(username, pw_hash, student_name, class_id, education_level=teacher_edu_level)
        students.append({"username": username, "name": student_name})

    return {
        "class_id": class_id,
        "class_code": class_code,
        "class_name": req.class_name,
        "password": req.password,
        "students": students,
    }


@app.get("/teacher/classes")
def list_teacher_classes(user: dict = Depends(get_current_user)):
    if user.get("role") != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can view classes")
    teacher_id = str(user["_id"])
    classes = get_classes_for_teacher(teacher_id)
    result = []
    for cls in classes:
        students = get_students_in_class(str(cls["_id"]))
        result.append({
            "class_id": str(cls["_id"]),
            "class_name": cls["class_name"],
            "class_code": cls["class_code"],
            "password": cls.get("password", ""),
            "student_count": cls["student_count"],
            "created_at": cls.get("created_at", ""),
            "settings": get_class_settings(cls),
            "students": [
                {"username": s.get("username"), "name": s.get("name")}
                for s in students
            ],
        })
    return {"classes": result}


class ClassSettingsReq(BaseModel):
    ai_enabled: Optional[bool] = None
    access_mode: Optional[str] = None
    unlocked_phase: Optional[int] = None


@app.patch("/teacher/class/{class_id}/settings")
def update_class_settings_endpoint(
    class_id: str,
    req: ClassSettingsReq,
    user: dict = Depends(get_current_user),
):
    """Teacher-controlled class modes (e.g. turn the AI assistant on/off)."""
    if user.get("role") != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can change class settings")
    cls = find_class_by_id(class_id)
    if not cls:
        raise HTTPException(status_code=404, detail="Class not found")
    if str(cls.get("teacher_id")) != str(user["_id"]):
        raise HTTPException(status_code=403, detail="You do not own this class")

    updates = {k: v for k, v in req.dict().items() if v is not None}
    if req.access_mode is not None and req.access_mode not in ("full", "step", "phase"):
        raise HTTPException(status_code=400, detail="access_mode must be full, step, or phase")
    settings = update_class_settings(class_id, updates)
    if settings is None:
        raise HTTPException(status_code=400, detail="No valid settings to update")
    return {"class_id": class_id, "settings": settings}


@app.get("/teacher/student-sessions")
def get_teacher_student_sessions(user: dict = Depends(get_current_user)):
    """Return all sessions for all students in the teacher's classes."""
    if user.get("role") != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can view student sessions")
    teacher_id = str(user["_id"])
    sessions = get_all_student_sessions_for_teacher(teacher_id)
    for s in sessions:
        s["_id"] = str(s["_id"])
        if "user" in s and "_id" in s["user"]:
            s["user"]["_id"] = str(s["user"]["_id"])
        s["completed_steps"] = _completed_steps(s.get("step_notes"), s.get("worldview_band"))
        # When the teacher last left feedback (drives the "awaiting your
        # feedback" dashboard tile)
        fb = s.get("teacher_feedback") or []
        s["last_feedback_at"] = max((f.get("created_at", "") for f in fb), default=None) or None
        # Remove bulky fields from the response
        s.pop("step_notes", None)
        s.pop("teacher_feedback", None)
        s.pop("chat", None)
    return {"sessions": sessions}


# ============================================================
# Teacher: View Student Design & Feedback
# ============================================================

def _verify_teacher_owns_student(session_id: str, teacher_user: dict):
    """Verify that the teacher owns the class the student (session owner) belongs to.
    Returns (session_doc, student_user) or raises 403."""
    from bson import ObjectId as _ObjId
    from database import classes_col as _classes_col

    doc = find_session(session_id)
    if not doc:
        raise HTTPException(status_code=404, detail="Session not found")
    student_user_id = doc.get("user_id")
    if not student_user_id:
        raise HTTPException(status_code=404, detail="Session has no user")
    student = find_user_by_id(student_user_id)
    if not student:
        raise HTTPException(status_code=404, detail="Student not found")
    student_class_id = student.get("class_id")
    if not student_class_id:
        raise HTTPException(status_code=403, detail="Student is not in a class")
    cls = _classes_col.find_one({"_id": _ObjId(student_class_id)})
    if not cls:
        raise HTTPException(status_code=404, detail="Class not found")
    if cls.get("teacher_id") != str(teacher_user["_id"]):
        raise HTTPException(status_code=403, detail="You do not own this student's class")
    return doc, student


@app.get("/teacher/student-session/{session_id}")
def get_teacher_student_session(session_id: str, user: dict = Depends(get_current_user)):
    """Return full session data for a student (teacher/admin view — excludes chat)."""
    if user.get("role") not in ("teacher", "admin"):
        raise HTTPException(status_code=403, detail="Only teachers/admins can view student sessions")
    if user.get("role") == "admin":
        doc = find_session(session_id)
        if not doc:
            raise HTTPException(status_code=404, detail="Session not found")
        student = find_user_by_id(doc.get("user_id", ""))
        if not student:
            raise HTTPException(status_code=404, detail="Student not found")
    else:
        doc, student = _verify_teacher_owns_student(session_id, user)

    step_notes = doc.get("step_notes") or {}
    completed = _compute_completed_steps_from_doc(doc)

    return {
        "session_id": doc["session_id"],
        "student_name": student.get("username") or student.get("name", ""),
        "student_username": student.get("username") or student.get("email") or "",
        "worldview_label": doc.get("worldview_label"),
        "resolved_path": doc.get("resolved_path"),
        "chosen_methodology": doc.get("chosen_methodology"),
        "active_step": doc.get("active_step", 1),
        "completed_steps": completed,
        "step_notes": step_notes,
        "teacher_feedback": doc.get("teacher_feedback", []),
    }


@app.get("/teacher/student-step-config")
def get_teacher_student_step_config(
    session_id: str = Query(...),
    step: int = Query(...),
    user: dict = Depends(get_current_user),
):
    """Return step config for a student's session (teacher/admin view — same logic as /step/config)."""
    if user.get("role") not in ("teacher", "admin"):
        raise HTTPException(status_code=403, detail="Only teachers/admins can view student step config")
    if user.get("role") == "teacher":
        _verify_teacher_owns_student(session_id, user)

    # Reuse the same config logic
    sess = _require_session(session_id)
    paths_cfg = load_paths_config()

    if step <= 3:
        return StepConfigResp(step=step)

    resolved = sess.resolved_path
    if not resolved:
        return StepConfigResp(step=step, title=f"Step {step}",
                              directions="Student has not selected a worldview yet.")

    all_paths = paths_cfg.get("paths", {})
    default_path = resolved  # the worldview's usual pathway
    # Non-mixed path where the student chose the other methodology at Step 4:
    # serve the chosen pathway's configuration (worldviews are defaults, not
    # locks - mirrors the LLM-guidance override in _step_llm_guidance)
    if (resolved != "mixed" and sess.chosen_methodology
            and sess.chosen_methodology != resolved
            and sess.chosen_methodology in all_paths):
        resolved = sess.chosen_methodology
    path_data = all_paths.get(resolved, {})
    step_key = str(step)
    step_cfg = path_data.get("steps", {}).get(step_key, {})

    if resolved == "mixed" and step >= 5 and step_cfg.get("inherits_from_chosen_methodology"):
        chosen = sess.chosen_methodology
        if not chosen:
            return StepConfigResp(step=step, path="mixed", title=f"Step {step}",
                                  directions="Student has not chosen a methodology yet.")
        inherited_cfg = all_paths.get(chosen, {}).get("steps", {}).get(step_key, {})
        guidance = inherited_cfg.get("llm_guidance", "")
        addendum = step_cfg.get("llm_guidance_addendum", "")
        if addendum:
            guidance = f"{guidance}\n{addendum}" if guidance else addendum
        return StepConfigResp(
            step=step, path="mixed",
            title=inherited_cfg.get("title", f"Step {step}"),
            directions=inherited_cfg.get("directions", ""),
            field_type=inherited_cfg.get("field_type"),
            field_key=inherited_cfg.get("field_key"),
            options=inherited_cfg.get("options"),
            fields=inherited_cfg.get("fields"),
            llm_guidance=guidance or None,
        )

    quant_opts = None
    qual_opts = None
    if resolved == "mixed" and step == 4:
        quant_opts = all_paths.get("quantitative", {}).get("steps", {}).get("4", {}).get("options")
        qual_opts = all_paths.get("qualitative", {}).get("steps", {}).get("4", {}).get("options")

    return StepConfigResp(
        step=step, path=resolved,
        default_path=default_path,
        title=step_cfg.get("title", f"Step {step}"),
        directions=step_cfg.get("directions", ""),
        field_type=step_cfg.get("field_type"),
        field_key=step_cfg.get("field_key"),
        options=step_cfg.get("options"),
        fields=step_cfg.get("fields"),
        llm_guidance=step_cfg.get("llm_guidance"),
        quantitative_options=quant_opts,
        qualitative_options=qual_opts,
    )


class TeacherFeedbackReq(BaseModel):
    session_id: str
    step: Optional[int] = None
    text: str


@app.post("/teacher/feedback")
def post_teacher_feedback(req: TeacherFeedbackReq, user: dict = Depends(get_current_user)):
    """Teacher/admin submits feedback for a student's session."""
    if user.get("role") not in ("teacher", "admin"):
        raise HTTPException(status_code=403, detail="Only teachers/admins can post feedback")
    if user.get("role") == "teacher":
        _verify_teacher_owns_student(req.session_id, user)

    feedback_item = {
        "id": str(uuid.uuid4()),
        "teacher_id": str(user["_id"]),
        "teacher_name": user.get("name", "Teacher"),
        "step": req.step,
        "text": req.text,
        "created_at": datetime.utcnow().isoformat() + "Z",
        "read": False,
    }

    from database import sessions_col
    sessions_col.update_one(
        {"session_id": req.session_id},
        {"$push": {"teacher_feedback": feedback_item}},
    )

    return {"ok": True, "feedback": feedback_item}


@app.get("/teacher/feedback/{session_id}")
def get_teacher_feedback(session_id: str, user: dict = Depends(get_current_user)):
    """Get all teacher feedback for a student's session."""
    if user.get("role") != "teacher":
        raise HTTPException(status_code=403, detail="Only teachers can view feedback")
    _verify_teacher_owns_student(session_id, user)
    doc = find_session(session_id)
    return {"feedback": doc.get("teacher_feedback", []) if doc else []}


# ---- Student feedback endpoints ----

@app.get("/student/feedback")
def get_student_feedback(user: dict = Depends(get_current_user)):
    """Return feedback for the current student's session."""
    user_id = str(user["_id"])
    from database import get_latest_session_for_user
    doc = get_latest_session_for_user(user_id)
    if not doc:
        return {"feedback": [], "unread_count": 0}
    feedback = doc.get("teacher_feedback", [])
    unread = sum(1 for f in feedback if not f.get("read"))
    return {"feedback": feedback, "unread_count": unread}


@app.post("/student/feedback/mark-read")
def mark_student_feedback_read(user: dict = Depends(get_current_user)):
    """Mark all feedback as read for the current student's session."""
    user_id = str(user["_id"])
    from database import get_latest_session_for_user, sessions_col
    doc = get_latest_session_for_user(user_id)
    if not doc:
        return {"ok": True}
    sessions_col.update_one(
        {"session_id": doc["session_id"], "teacher_feedback": {"$exists": True}},
        {"$set": {"teacher_feedback.$[elem].read": True}},
        array_filters=[{"elem.read": False}],
    )
    return {"ok": True}


# ============================================================
# Routes
# ============================================================
@app.post("/session", response_model=SessionCreateResponse)
def create_session(user: dict = Depends(get_current_user)):
    sid = str(uuid.uuid4())
    user_id = str(user["_id"])
    create_session_doc(sid, user_id)
    return SessionCreateResponse(session_id=sid)


class SessionResumeResponse(BaseModel):
    session_id: Optional[str] = None
    active_step: int = 1
    found: bool = False
    completed_steps: List[int] = []


def _step_has_content(data: dict) -> bool:
    """True when a step's saved data holds at least one non-empty value -
    a dict full of empty strings (created by typing then deleting) does NOT
    count as a completed step."""
    for v in (data or {}).values():
        if isinstance(v, str):
            if v.strip():
                return True
        elif v:  # non-empty list/number/bool/dict
            return True
    return False


def _completed_steps(notes: dict, worldview_band: Optional[str] = None) -> List[int]:
    """SINGLE source of truth for step completion (used by every endpoint).

    Step 1 is complete when a worldview was chosen - via the step notes OR the
    session-level worldview_band, so a later /step/save that overwrites the
    notes can never un-complete it. Steps 2-9 are complete when they hold at
    least one non-empty value."""
    notes = notes or {}
    completed = []
    s1 = notes.get("1") or {}
    if s1.get("worldview_id") or s1.get("worldview") or worldview_band:
        completed.append(1)
    for s in range(2, 10):
        if _step_has_content(notes.get(str(s)) or {}):
            completed.append(s)
    return completed


def _compute_completed_steps_from_session(sess: "SessionData") -> List[int]:
    return _completed_steps(sess.step_notes, sess.worldview_band)


def _compute_completed_steps_from_doc(doc: dict) -> List[int]:
    return _completed_steps(doc.get("step_notes"), doc.get("worldview_band"))


@app.get("/session/resume", response_model=SessionResumeResponse)
def resume_session(user: dict = Depends(get_current_user)):
    """Return the user's most recent session if one exists."""
    user_id = str(user["_id"])
    doc = get_latest_session_for_user(user_id)
    if not doc:
        return SessionResumeResponse(found=False)
    return SessionResumeResponse(
        session_id=doc["session_id"],
        active_step=doc.get("active_step", 1),
        found=True,
        completed_steps=_compute_completed_steps_from_doc(doc),
    )


class SessionSummary(BaseModel):
    session_id: str
    created_at: str = ""
    active_step: int = 1
    completed_steps: List[int] = []
    topic: Optional[str] = None
    resolved_path: Optional[str] = None
    worldview_label: Optional[str] = None


class SessionListResponse(BaseModel):
    sessions: List[SessionSummary]


@app.get("/session/list", response_model=SessionListResponse)
def list_sessions(user: dict = Depends(get_current_user)):
    """Return all sessions for the current user with summary info."""
    user_id = str(user["_id"])
    docs = get_session_summaries_for_user(user_id)
    summaries = []
    for doc in docs:
        step_notes = doc.get("step_notes") or {}
        topic = (step_notes.get("2") or {}).get("topic")
        summaries.append(SessionSummary(
            session_id=doc["session_id"],
            created_at=doc.get("created_at", ""),
            active_step=doc.get("active_step", 1),
            completed_steps=_compute_completed_steps_from_doc(doc),
            topic=topic,
            # Display the pathway the student actually chose (e.g. a
            # constructivist running a quantitative study), falling back to
            # the worldview's resolved default
            resolved_path=(
                doc.get("chosen_methodology")
                if doc.get("chosen_methodology") and doc.get("resolved_path") != "mixed"
                else doc.get("resolved_path")
            ),
            worldview_label=doc.get("worldview_label"),
        ))
    return SessionListResponse(sessions=summaries)


class UpdateStepReq(BaseModel):
    session_id: str
    active_step: int


@app.post("/session/update_step")
def update_active_step(req: UpdateStepReq, user: dict = Depends(get_current_user)):
    """Save the student's current active step to the session."""
    sess = _require_session(req.session_id)
    _guard_step_access(user, sess, req.active_step)
    sess.active_step = req.active_step
    _persist_session(sess)
    return {"ok": True}


@app.get("/chat/history", response_model=ChatHistoryResp)
def get_chat_history(session_id: str = Query(...), user: dict = Depends(get_current_user)):
    sess = _require_session(session_id)
    history = _get_chat(sess)
    return ChatHistoryResp(session_id=session_id, history=history)

class StepDataReq(BaseModel):
    session_id: str
    step: int
    data: Dict[str, Any]


class StepDataResp(BaseModel):
    session_id: str
    step: int
    data: Dict[str, Any]
    completed_steps: List[int] = []


@app.post("/step/save", response_model=StepDataResp)
def save_step_data(req: StepDataReq, user: dict = Depends(get_current_user)):
    sess = _require_session(req.session_id)
    _guard_step_access(user, sess, req.step)
    key = str(req.step)
    # MERGE into the existing notes instead of replacing: the frontend only
    # knows its own form fields, so a replace used to wipe backend-written
    # keys (worldview_id from /worldview/set, chosen_methodology from
    # /step/set_methodology), silently un-completing steps.
    sess.step_notes[key] = {**(sess.step_notes.get(key) or {}), **(req.data or {})}
    _persist_session(sess)
    return StepDataResp(session_id=sess.id, step=req.step, data=sess.step_notes[key], completed_steps=_compute_completed_steps_from_session(sess))


@app.get("/step/get", response_model=StepDataResp)
def get_step_data(
    session_id: str = Query(...),
    step: int = Query(...),
    user: dict = Depends(get_current_user),
):
    sess = _require_session(session_id)
    key = str(step)
    data = sess.step_notes.get(key, {})
    return StepDataResp(session_id=sess.id, step=step, data=data)

class WorldviewSetReq(BaseModel):
    session_id: str
    worldview_id: str  # "positivist" | "post_positivist" | "constructivist" | "transformative" | "pragmatist"

class WorldviewSetResp(BaseModel):
    session_id: str
    worldview_id: str
    worldview_label: str
    completed_steps: List[int] = []


WORLDVIEW_LABELS = {
    "positivist": "Positivist",
    "post_positivist": "Post Positivist",
    "constructivist": "Constructivist",
    "transformative": "Transformative",
    "pragmatist": "Pragmatist",
}


@app.post("/worldview/set", response_model=WorldviewSetResp)
def set_worldview(req: WorldviewSetReq, user: dict = Depends(get_current_user)):
    sess = _require_session(req.session_id)

    wid = (req.worldview_id or "").strip()
    if wid not in WORLDVIEW_LABELS:
        raise HTTPException(status_code=400, detail="Invalid worldview_id")

    # Set worldview on session (what the LLM will see)
    sess.worldview_band = wid
    sess.worldview_label = WORLDVIEW_LABELS[wid]

    # Also store into step 1 notes so you can use it later
    # Save both "worldview_id" (used by _compute_completed_steps) and
    # "worldview" (used by the frontend dropdown to display selection)
    sess.step_notes["1"] = {**(sess.step_notes.get("1") or {}), "worldview_id": wid, "worldview": wid}

    # Resolve research path from worldview
    paths_cfg = load_paths_config()
    wv_to_path = paths_cfg.get("worldview_to_path", {})
    sess.resolved_path = wv_to_path.get(wid, None)

    _persist_session(sess)
    return WorldviewSetResp(
        session_id=sess.id,
        worldview_id=wid,
        worldview_label=sess.worldview_label,
        completed_steps=_compute_completed_steps_from_session(sess),
    )

# ---------------- Step config + methodology endpoints ----------------

class StepConfigResp(BaseModel):
    step: int
    path: Optional[str] = None
    title: str = ""
    directions: str = ""
    field_type: Optional[str] = None
    field_key: Optional[str] = None
    options: Optional[List[Dict[str, Any]]] = None
    fields: Optional[List[Dict[str, Any]]] = None
    llm_guidance: Optional[str] = None
    quantitative_options: Optional[List[Dict[str, Any]]] = None
    qualitative_options: Optional[List[Dict[str, Any]]] = None
    recommended_methodology: Optional[str] = None  # "quantitative" or "qualitative"
    default_path: Optional[str] = None  # the worldview's usual pathway (path may differ if the student switched)


@app.get("/step/config", response_model=StepConfigResp)
def get_step_config(
    session_id: str = Query(...),
    step: int = Query(...),
    lang: Optional[str] = Query(None),
    user: dict = Depends(get_current_user),
):
    """Return the path-resolved configuration for a given step."""
    sess = _require_session(session_id)
    paths_cfg = load_paths_config()
    # The client sends its current language explicitly; the account setting is
    # only a fallback (it may lag behind a just-made switch in the UI).
    lang = lang or (user or {}).get("language", "en")

    # Steps 1-3 are handled entirely by the frontend
    if step <= 3:
        return StepConfigResp(step=step)

    resolved = sess.resolved_path
    if not resolved:
        return StepConfigResp(
            step=step,
            title=f"Step {step}",
            directions="Please complete Step 1 first and select your worldview.",
        )

    all_paths = paths_cfg.get("paths", {})
    path_data = all_paths.get(resolved, {})
    step_key = str(step)
    step_cfg = path_data.get("steps", {}).get(step_key, {})

    # --- Steps 5-9: resolve from chosen_methodology if student overrode ---
    # Works for mixed (inherits_from_chosen_methodology) AND for non-mixed overrides
    effective_path = sess.chosen_methodology or resolved
    if step >= 5:
        # Mixed paths with explicit inheritance flag
        if resolved == "mixed" and step_cfg.get("inherits_from_chosen_methodology"):
            chosen = sess.chosen_methodology
            if not chosen:
                return StepConfigResp(
                    step=step,
                    path="mixed",
                    title=f"Step {step}",
                    directions="Please complete Step 4 first and choose your primary methodology.",
                )
            inherited_cfg = all_paths.get(chosen, {}).get("steps", {}).get(step_key, {})
            guidance = inherited_cfg.get("llm_guidance", "")
            inherited_cfg = _localize_step_cfg(chosen, step_key, inherited_cfg, lang)
            addendum = step_cfg.get("llm_guidance_addendum", "")
            if addendum:
                guidance = f"{guidance}\n{addendum}" if guidance else addendum
            return StepConfigResp(
                step=step,
                path=resolved,
                title=inherited_cfg.get("title", f"Step {step}"),
                directions=inherited_cfg.get("directions", ""),
                field_type=inherited_cfg.get("field_type"),
                field_key=inherited_cfg.get("field_key"),
                options=inherited_cfg.get("options"),
                fields=inherited_cfg.get("fields"),
                llm_guidance=guidance or None,
            )
        # Non-mixed path but student overrode methodology at Step 4
        if resolved != "mixed" and sess.chosen_methodology and sess.chosen_methodology != resolved:
            override_cfg = all_paths.get(sess.chosen_methodology, {}).get("steps", {}).get(step_key, {})
            override_cfg = _localize_step_cfg(sess.chosen_methodology, step_key, override_cfg, lang)
            if override_cfg:
                return StepConfigResp(
                    step=step,
                    path=resolved,
                    title=override_cfg.get("title", f"Step {step}"),
                    directions=override_cfg.get("directions", ""),
                    field_type=override_cfg.get("field_type"),
                    field_key=override_cfg.get("field_key"),
                    options=override_cfg.get("options"),
                    fields=override_cfg.get("fields"),
                    llm_guidance=override_cfg.get("llm_guidance"),
                )

    # --- Step 4: methodology decision for ALL paths ---
    if step == 4:
        step_cfg = _localize_step_cfg(resolved, "4", step_cfg, lang)
        quant_opts = (
            _localize_step_cfg("quantitative", "4",
                               all_paths.get("quantitative", {}).get("steps", {}).get("4", {}), lang)
            .get("options")
        )
        qual_opts = (
            _localize_step_cfg("qualitative", "4",
                               all_paths.get("qualitative", {}).get("steps", {}).get("4", {}), lang)
            .get("options")
        )
        # Determine recommendation based on worldview path
        recommended = None
        if resolved == "quantitative":
            recommended = "quantitative"
        elif resolved == "qualitative":
            recommended = "qualitative"
        # For mixed (pragmatist): no recommendation — both are equally valid

        # Once the student has confirmed a methodology, the title/directions
        # follow their CHOICE, not the worldview's default pathway (a
        # constructivist doing a quantitative study must not see
        # "(Qualitative Design)" in the header).
        title_cfg = step_cfg
        if resolved != "mixed" and sess.chosen_methodology and sess.chosen_methodology in all_paths:
            title_cfg = _localize_step_cfg(
                sess.chosen_methodology, "4",
                all_paths[sess.chosen_methodology].get("steps", {}).get("4", step_cfg), lang)
        return StepConfigResp(
            step=step,
            path=resolved,
            title=title_cfg.get("title", f"Step 4: How will I study it?"),
            directions=title_cfg.get("directions", ""),
            field_type="methodology_decision",
            field_key=step_cfg.get("field_key", "design"),
            options=step_cfg.get("options"),
            fields=step_cfg.get("fields"),
            llm_guidance=step_cfg.get("llm_guidance"),
            quantitative_options=quant_opts,
            qualitative_options=qual_opts,
            recommended_methodology=recommended,
        )

    step_cfg = _localize_step_cfg(resolved, step_key, step_cfg, lang)
    return StepConfigResp(
        step=step,
        path=resolved,
        title=step_cfg.get("title", f"Step {step}"),
        directions=step_cfg.get("directions", ""),
        field_type=step_cfg.get("field_type"),
        field_key=step_cfg.get("field_key"),
        options=step_cfg.get("options"),
        fields=step_cfg.get("fields"),
        llm_guidance=step_cfg.get("llm_guidance"),
    )


class SetMethodologyReq(BaseModel):
    session_id: str
    methodology: str  # "quantitative" or "qualitative"


class SetMethodologyResp(BaseModel):
    session_id: str
    chosen_methodology: str


@app.post("/step/set_methodology", response_model=SetMethodologyResp)
def set_methodology(req: SetMethodologyReq, user: dict = Depends(get_current_user)):
    """Set the primary methodology chosen at Step 4 (all paths, not just mixed)."""
    sess = _require_session(req.session_id)
    meth = (req.methodology or "").strip().lower()
    if meth not in ("quantitative", "qualitative"):
        raise HTTPException(
            status_code=400,
            detail="methodology must be 'quantitative' or 'qualitative'",
        )

    # If changing pathways, clear the now-mismatched later-step data. For
    # non-mixed paths the baseline is the worldview's resolved path, so the
    # FIRST switch away from the default also clears (the saved Step 4 design
    # and steps 5-9 belong to the other pathway).
    prev = sess.chosen_methodology
    prev_effective = prev or (sess.resolved_path if sess.resolved_path != "mixed" else None)
    if prev_effective and prev_effective != meth:
        for s in range(5, 10):
            sess.step_notes.pop(str(s), None)
        if sess.resolved_path != "mixed":
            s4_old = sess.step_notes.get("4") or {}
            s4_old.pop("design", None)
            sess.step_notes["4"] = s4_old

    sess.chosen_methodology = meth
    s4 = sess.step_notes.get("4") or {}
    s4["chosen_methodology"] = meth
    sess.step_notes["4"] = s4

    _persist_session(sess)
    return SetMethodologyResp(session_id=sess.id, chosen_methodology=meth)



def _render_step_context(sess: SessionData) -> str:
    """Build a comprehensive context string from all step inputs (1-9)."""
    lines = []
    notes = sess.step_notes or {}

    # Step 1: worldview (field saved as "worldview_id" by /worldview/set)
    s1 = notes.get("1") or {}
    worldview_id = (s1.get("worldview_id") or s1.get("worldview") or "").strip()
    if worldview_id:
        label = WORLDVIEW_LABELS.get(worldview_id, worldview_id)
        lines.append(f"Step 1 worldview: {label}")

    # Step 2: topic and goals (personal, practical, intellectual)
    s2 = notes.get("2") or {}
    if s2.get("topic"):
        lines.append(f"Step 2 research topic: {s2['topic']}")
    if s2.get("personalGoals"):
        lines.append(f"Step 2 personal goals: {s2['personalGoals']}")
    if s2.get("practicalGoals"):
        lines.append(f"Step 2 practical goals: {s2['practicalGoals']}")
    if s2.get("intellectualGoals"):
        lines.append(f"Step 2 intellectual goals: {s2['intellectualGoals']}")
    # backward compat: old single "goals" field
    if s2.get("goals") and not s2.get("personalGoals"):
        lines.append(f"Step 2 research goals: {s2['goals']}")

    # Step 3: literature review
    s3 = notes.get("3") or {}
    if s3.get("topicalResearch"):
        lines.append(f"Step 3 topical research: {s3['topicalResearch']}")
    if s3.get("theoreticalFrameworks"):
        lines.append(f"Step 3 theoretical frameworks: {s3['theoreticalFrameworks']}")

    # Resolved path and methodology (if set)
    if sess.resolved_path:
        lines.append(f"Research path: {sess.resolved_path}")
    if sess.chosen_methodology:
        lines.append(f"Chosen methodology (Step 4): {sess.chosen_methodology}")

    # Steps 4-9: read whatever structured data was saved
    for step_num in range(4, 10):
        sn = notes.get(str(step_num)) or {}
        for key, val in sn.items():
            if not val:
                continue
            if isinstance(val, list):
                val_str = ", ".join(str(v) for v in val)
            else:
                val_str = str(val)
            # Truncate very long values to keep context manageable
            if len(val_str) > 300:
                val_str = val_str[:300] + "..."
            lines.append(f"Step {step_num} {key}: {val_str}")

    return "\n".join(lines) if lines else "No step inputs saved yet."


# ---------------- Chat main endpoint (non-streaming) ----------------
# ---------------- Teacher-controlled AI mode ----------------

_AI_OFF_MESSAGE = (
    "The AI assistant is currently turned off for your class. Your teacher has "
    "chosen to have you work through your research design on your own for now. "
    "Keep going in the 'My Research Design' panel — and check back later, as your "
    "teacher may turn the assistant on when you're ready."
)


# ---------------- Harmful-content safety guardrail (Llama Guard 3) ----------------

_SAFETY_REFUSAL = (
    "I can't help with that request. Hopscotch is a research-methods learning space, "
    "and I'm here to support safe, ethical, and academic work. If this connects to a "
    "genuine research interest, try reframing it around how you would study the topic "
    "responsibly and ethically — and I'll gladly help you design that study."
)

_AI_OFF_MESSAGE_ES = (
    "Tu docente ha desactivado el asistente de IA para tu clase por ahora. "
    "Puedes seguir trabajando en cada paso de tu diseño de investigación; "
    "tus respuestas se guardan normalmente."
)

_SAFETY_REFUSAL_ES = (
    "No puedo ayudarte con esa solicitud. Hopscotch es un espacio de aprendizaje "
    "de métodos de investigación, y estoy aquí para apoyar un trabajo seguro, ético "
    "y académico. Si esto se conecta con un interés de investigación genuino, "
    "intenta replantearlo en torno a cómo estudiarías el tema de forma responsable "
    "y ética — y con gusto te ayudaré a diseñar ese estudio."
)

_AI_OFF_MESSAGE_ZH = (
    "你的老师目前为班级关闭了 AI 助手，希望大家先自己完成研究设计。"
    "你可以继续在“我的研究设计”面板中作答，内容会正常保存——"
    "之后老师可能会在合适的时候重新开启助手。"
)

_SAFETY_REFUSAL_ZH = (
    "我无法帮助处理这个请求。Hopscotch 是一个学习研究方法的空间，"
    "我在这里支持安全、合乎伦理的学术工作。如果这与你真实的研究兴趣有关，"
    "请试着换个角度：思考如何以负责任、合乎伦理的方式研究这个主题——"
    "我很乐意帮你设计这样的研究。"
)


def _canned(user: dict, en_text: str, es_text: str, zh_text: str = None) -> str:
    lang = (user or {}).get("language")
    if lang == "es":
        return es_text
    if lang == "zh" and zh_text:
        return zh_text
    return en_text



def _moderate_input(user_msg: str) -> tuple:
    """Run Llama Guard 3 locally (via Ollama) on the student's message.
    Returns (is_safe: bool, categories: str). Fails OPEN (treated as safe) if the
    guard model is unavailable, so moderation can never break the chat."""
    if not MODERATION_ENABLED or not (user_msg or "").strip():
        return True, ""
    try:
        r = requests.post(OLLAMA_URL, json={
            "model": MODERATION_MODEL,
            "messages": [{"role": "user", "content": user_msg}],
            "stream": False,
            "options": {"temperature": 0.0, "num_predict": 20},
        }, timeout=20)
        r.raise_for_status()
        out = ((r.json().get("message", {}) or {}).get("content", "") or "").strip()
    except Exception as e:
        logger.warning("Moderation (Llama Guard) unavailable — allowing message: %s", e)
        return True, ""
    lines = [l.strip() for l in out.splitlines() if l.strip()]
    if lines and lines[0].lower().startswith("unsafe"):
        categories = lines[1] if len(lines) > 1 else ""
        logger.info("Llama Guard flagged message as unsafe (%s)", categories or "unspecified")
        return False, categories
    return True, ""


# ---------------- Academic-integrity guardrail (coach, don't author) ----------------

# Step-specific guiding nudges used when redirecting a "do it for me" request.
_COACH_NUDGES = {
    1: "What draws you toward a particular worldview? Consider whether you see knowledge as objective and measurable, or as socially constructed and shaped by context.",
    2: "What issue in your field genuinely interests you? What problem have you noticed that you'd like to understand or change?",
    3: "What research have you already read on this topic? Which theories from your coursework feel connected to it?",
    4: "Given your worldview, are you trying to measure something or understand people's experiences? What kind of data would best answer that?",
    5: "What exactly are you trying to find out? Try turning your topic into a question (qualitative) or a testable prediction/hypothesis (quantitative).",
    6: "Who could give you the information you need, and how might you gather it — survey, interview, observation, existing records?",
    7: "Once you have your data, what would you need to do to make sense of it and answer your research question?",
    8: "What could make someone doubt your findings, and what steps would reduce that doubt?",
    9: "Who might your study affect, and how will you protect participants and act responsibly?",
}


# ---------------- Source/citation guardrail (deterministic, model-proof) ----------------
# The LLM has been observed producing (and fabricating) specific citations despite the
# system prompt forbidding it. This regex gate runs BEFORE the model, so it cannot be
# overridden. It blocks requests that ask the AI to PRODUCE sources, while still allowing
# "how/where do I find sources?" and "how do I cite?" questions.

_SOURCE_TARGET = (
    r"(citations?|references?|articles?|papers?|studies|source list|sources?|"
    r"journal articles?|bibliograph\w*|works cited|literature list)"
)
_PRODUCE_VERB = (
    r"(give|provide|list|share|write|generate|recommend|suggest|show|send|name|"
    r"find\s+me|get\s+me)\s+(me\s+|us\s+)?(some\s+|a\s+few\s+|the\s+|specific\s+|"
    r"actual\s+|real\s+|relevant\s+|several\s+|examples?\s+of\s+)?"
)


def _asks_for_sources(user_msg: str) -> bool:
    """True if the student is asking the AI to PRODUCE citations/references/sources
    (which risks fabrication and does their literature work for them). Deliberately
    permits 'where/how can I find…' and 'how do I cite…' questions."""
    m = (user_msg or "").lower()
    if not m:
        return False
    # Allow genuine "how to find / how to cite / how to format" questions.
    if re.search(r"how\s+(do|can|should|would)\s+i\s+(cite|format|reference)\b", m):
        return False
    # Block: "the citations/sources you recommended/gave/mentioned/listed"
    if re.search(_SOURCE_TARGET + r"\s+(you|u)\s+(recommend|mention|gave|give|suggest|provid|list|cite)\w*", m):
        return True
    # Block: "title(s) of the article(s)/paper(s)/study/studies/source(s)"
    if re.search(r"titles?\s+(of|for)\s+(the\s+|those\s+|these\s+)?(articles?|papers?|studies|study|sources?|references?)", m):
        return True
    # Block: a 'produce' verb followed by a source target ("give me references",
    # "provide sources", "recommend studies on X", "find me articles").
    if re.search(_PRODUCE_VERB + _SOURCE_TARGET, m):
        return True
    # Block: question forms that ask the AI to surface sources
    # ("what are some studies on X", "any articles about Y", "which papers discuss Z").
    if re.search(r"\b(what|which|are there any|any|got any|know any)\b[^.?!]{0,22}\b"
                 + _SOURCE_TARGET
                 + r"\b[^.?!]{0,18}\b(on|about|for|regarding|related|discuss|examin|support|show|cover)", m):
        return True
    return False


def _source_redirect_message(active_step) -> str:
    """Returned instead of letting the model invent or hand over citations."""
    return (
        "I'm not going to hand you specific citations or article titles — for two "
        "reasons. First, finding and judging sources is a core research skill this step "
        "is here to build. Second, an AI can get references wrong or invent ones that "
        "don't exist, and citing a fake source would seriously hurt your work.\n\n"
        "Here's how to find strong, real sources yourself:\n"
        "• Search Google Scholar (scholar.google.com) or your library databases "
        "(PsycINFO, PubMed, ERIC, JSTOR) using 2-3 keywords from your topic.\n"
        "• Skim the abstract to check it actually fits your question, and follow the "
        "reference lists of good papers to find more.\n\n"
        "Once you've found a study, paste what it says or what it found, and I'll help "
        "you analyze how it connects to your topic and where the gaps are."
    )


# ---------------- Output guard: strip handed-over deliverables ----------------
# Even with a strong system prompt, the local model sometimes hands the student a
# finished "Revised Topic: …" / "Example Research Question: …" block — i.e. it does
# their work. This deterministic filter removes those blocks from the output (both
# streaming and non-streaming) so the answer keeps the coaching but not the answer.

# Words the model uses to head a block where it hands the student a finished
# deliverable (their own design content). Matched only when the term is the WHOLE
# heading line (ends in ':' or end-of-line) — so it fires on "#### Personal Goals:"
# but not on the term used inside a normal feedback sentence.
_HANDED_ADJ = (
    r'(?:revised|refined|suggested|improved|reworded|polished|reformulated|'
    r'proposed|updated|possible|recommended|stronger|better|alternative|rewritten)'
)
# Deliverable field names (incl. short forms). Only ever matched WITH an adjective in
# front or a "suggestion/version" word after — never bare, so the model restating the
# student's own goals/topic for feedback is NOT stripped (that would fragment feedback).
_HANDED_DELIVERABLE = (
    r'(?:personal|practical|intellectual|research)\s+goals?'
    r'|research\s+topic|research\s+questions?|research\s+design|research\s+aims?'
    r'|problem\s+statement|topical\s+research|literature\s+review'
    r'|(?:theoretical|conceptual)\s+frameworks?'
    r'|data\s+(?:collection|gathering|analysis)|analysis\s+plan|hypothesis'
    r'|topic|questions?|goals?|aims?|objectives?'
)
_HANDED_SUGG = r'(?:suggestions?|rewrites?|revisions?|versions?|drafts?|examples?|rewrite)'
# Forward-preview headings ("Next Steps", "Moving Forward") — previewing later steps.
_HANDED_FWD = (
    r'(?:next\s+steps?|moving\s+forward|what\'?s\s+next|going\s+forward|'
    r'recommended\s+next\s+steps?|where\s+to\s+go\s+from\s+here|looking\s+ahead)'
)
_HANDED_LABEL_RE = re.compile(
    r'^[#>*_\s\d.\-]*(?:'
    # A: adjective + deliverable  ("Refined Topic", "Suggested Goals", "Refined Topic Suggestion")
    r'(?:' + _HANDED_ADJ + r')\s+(?:research\s+)?(?:' + _HANDED_DELIVERABLE + r')(?:\s+' + _HANDED_SUGG + r')?'
    # B: deliverable + suggestion word  ("Topic Suggestion", "Goals Suggestion")
    r'|(?:' + _HANDED_DELIVERABLE + r')\s+' + _HANDED_SUGG +
    # C: forward-preview section headings
    r'|' + _HANDED_FWD +
    r')\b[^A-Za-z]*(:|$)',
    re.I,
)
_HANDED_NUDGE = (
    "_(I've held back the exact wording here — that part is yours to write. Draft it in "
    "the 'My Research Design' panel and I'll give you feedback on what you come up with.)_"
)

# The model sometimes closes a message with a letter-style sign-off using a placeholder
# ("Warm regards, [Your Tutor]"). Normalize any tutor/name sign-off to "Hopscotch".
_SIGNOFF_PLACEHOLDER_RE = re.compile(
    r'\[[^\]\n]*?(?:tutor|your name|your ai|assistant)[^\]\n]*?\]', re.I)
_SIGNOFF_LINE_RE = re.compile(
    r'(?im)^(\s*(?:warm(?:est)?\s+regards|best(?:\s+regards|\s+wishes)?|kind\s+regards|'
    r'sincerely|warmly|cheers|regards|thanks|thank\s+you)\s*[,\-–—:]?\s*)'
    r'(?:your\s+)?(?:ai\s+)?(?:research\s+)?tutor\b\.?\s*$')


def _fix_signoff(s: str) -> str:
    """Replace placeholder / generic-tutor sign-offs with 'Hopscotch'."""
    if not s:
        return s
    s = _SIGNOFF_PLACEHOLDER_RE.sub("Hopscotch", s)
    s = _SIGNOFF_LINE_RE.sub(lambda m: m.group(1) + "Hopscotch", s)
    return s


def _sanitize_stream(raw_iter):
    """Output guard: drops handed-over deliverable blocks and normalizes
    AI-looking punctuation (em/en dashes) in everything the student sees."""
    for piece in _sanitize_stream_blocks(raw_iter):
        yield piece.replace("—", "-").replace("–", "-")


def _sanitize_stream_blocks(raw_iter):
    """Line-buffer a token stream and drop 'handed-over deliverable' blocks, emitting a
    single coaching nudge in their place. Yields sanitized text pieces."""
    pending = ""
    st = {"suppress": False, "seen_content": False, "nudges": 0}

    def start_block():
        # Drop held-back blocks SILENTLY — the student shouldn't see an internal
        # "I've held back the wording" note; they just get clean feedback.
        st["suppress"] = True
        st["seen_content"] = False
        return ""

    def process_line(line):
        s = line.strip()
        is_label = bool(_HANDED_LABEL_RE.match(s))
        if st["suppress"]:
            if is_label or s.startswith("#"):
                # A new label/heading ends the current block; handle this line fresh.
                st["suppress"] = False
            elif s == "":
                if not st["seen_content"]:
                    return ""  # blank sitting between the label and its value → drop
                st["suppress"] = False
                return "\n"       # blank after the value → block ends here
            else:
                st["seen_content"] = True
                return ""         # the handed-over value itself → drop
        if is_label:
            return start_block()
        return _fix_signoff(line) + "\n"

    for delta in raw_iter:
        pending += delta
        while "\n" in pending:
            line, pending = pending.split("\n", 1)
            out = process_line(line)
            if out:
                yield out
    if pending:
        out = process_line(pending)
        if out:
            yield out.rstrip("\n") if not pending.endswith("\n") else out


def _strip_handed_answers(text: str) -> str:
    """Non-streaming version of the output guard."""
    if not text:
        return text
    return "".join(_sanitize_stream(iter([text]))).rstrip("\n")


def _coach_redirect_message(active_step) -> str:
    """The warm 'I'll guide you, but you draft it' response returned instead of
    authoring the student's design content for them."""
    nudge = _COACH_NUDGES.get(active_step or 0,
        "Jot down your initial thoughts in the 'My Research Design' panel first.")
    return (
        "I'm here to coach you through this, not to write it for you, because doing it "
        "yourself is how you truly learn to design research. Let's build it together.\n\n"
        f"To get you started: {nudge}\n\n"
        "Write your first attempt in the 'My Research Design' panel and I'll give you "
        "detailed feedback and help you strengthen it."
    )


def _asks_ai_to_author(user_msg: str, active_step) -> bool:
    """Fast local self-check (intent gate): is the student asking the AI to PRODUCE
    their own design content (topic, research question, aim, hypothesis, or any step
    plan) rather than asking for an explanation, feedback, or guidance? Returns True
    only for 'do it for me' requests. Fails OPEN (False) if the classifier is
    unavailable — the strengthened system prompt still applies as a backstop."""
    prompt = (
        "You are a classifier for a research-methods tutoring app. A student is on "
        f"Step {active_step or '?'} of designing their OWN research study.\n"
        "Classify what their message is asking for. Reply with EXACTLY one word:\n"
        "AUTHOR = they want the AI to write/produce/generate/REWRITE THEIR OWN design "
        "content for them, or to hand them a finished version. Examples: 'write my "
        "research question', 'give me a hypothesis', 'what should my topic be', 'make my "
        "data-collection plan', 'just do it for me', 'give me the answer', 'write it for "
        "me', 'rewrite/reword/rephrase my question into a better version', 'give me a "
        "revised/improved version of my topic', 'give me citations/references/articles/"
        "sources', 'what are some studies on X'.\n"
        "AUTHOR also includes asking you to produce, draft, decide, or work out THEIR "
        f"content for a step LATER than Step {active_step or '?'} — e.g. (while on an "
        "earlier step) 'what methodology should I use', 'give me my research question', "
        "'how should I collect/analyze my data', 'what should my hypothesis be'. Producing "
        "their answers for steps they haven't reached counts as AUTHOR.\n"
        "COACH = they want a concept explained in general, a definition, feedback/critique "
        "on something they ALREADY wrote for the CURRENT step (strengths/weaknesses without "
        "rewriting it), an example that teaches an idea, guidance on HOW to find sources, or "
        "general guidance about the current step.\n\n"
        f"Student message: \"{(user_msg or '')[:500]}\"\n\n"
        "One word (AUTHOR or COACH):"
    )
    msgs = [{"role": "user", "content": prompt}]
    raw = None
    try:
        if LLM_BACKEND == "vllm":
            raw = _call_vllm(msgs, temperature=0.0, max_tokens=4, timeout=20)
        if raw is None:
            raw = _call_ollama({
                "model": LLM_MODEL, "messages": msgs, "stream": False,
                "options": {"temperature": 0.0, "num_predict": 4},
            }, timeout=20)
    except Exception as e:
        logger.warning("Intent gate classifier failed: %s", e)
        return False
    if not raw:
        return False
    # Take the first word to avoid stray tokens flipping the result
    return raw.strip().upper().startswith("AUTHOR") or "AUTHOR" in raw.strip().upper()[:12]


@app.post("/chat/send", response_model=ChatHistoryResp)
def chat_send(req: ChatSendReq = Body(...), user: dict = Depends(get_current_user)):
    sess = _require_session(req.session_id)
    history = _get_chat(sess)

    user_msg = (req.message or "").strip()
    if not user_msg:
        return ChatHistoryResp(session_id=req.session_id, history=history)

    # store user turn
    history.append(ChatTurn(role="user", content=user_msg, step=req.active_step))

    # Teacher-controlled mode: AI assistant may be turned off for this student's class.
    if not _student_ai_enabled(user):
        history.append(ChatTurn(role="assistant", content=_canned(user, _AI_OFF_MESSAGE, _AI_OFF_MESSAGE_ES, _AI_OFF_MESSAGE_ZH), step=req.active_step))
        _persist_session(sess)
        return ChatHistoryResp(session_id=req.session_id, history=history)

    # Source gate (deterministic): a citation/source request is never a safety concern,
    # so handle it before moderation to return the helpful coaching message.
    if _asks_for_sources(user_msg):
        answer = _source_redirect_message(req.active_step)
        history.append(ChatTurn(role="assistant", content=answer, step=req.active_step))
        _persist_session(sess)
        return ChatHistoryResp(session_id=req.session_id, history=history)

    # Safety gate: refuse harmful/unethical requests (Llama Guard 3, local).
    is_safe, _cats = _moderate_input(user_msg)
    if not is_safe:
        history.append(ChatTurn(role="assistant", content=_canned(user, _SAFETY_REFUSAL, _SAFETY_REFUSAL_ES, _SAFETY_REFUSAL_ZH), step=req.active_step))
        _persist_session(sess)
        return ChatHistoryResp(session_id=req.session_id, history=history)

    # Academic-integrity gate: if the student is asking the AI to author their design
    # content, coach them instead of doing the work for them.
    if _asks_ai_to_author(user_msg, req.active_step):
        answer = _coach_redirect_message(req.active_step)
        history.append(ChatTurn(role="assistant", content=answer, step=req.active_step))
        _persist_session(sess)
        return ChatHistoryResp(session_id=req.session_id, history=history)

    # Normal LLM + RAG chat using worldview and resources
    worldview_profile = _render_worldview_profile(sess)
    step_context = _render_step_context(sess)
    passages = _retrieve(user_msg, k=5)
    step_llm_guidance = _get_step_llm_guidance(sess, req.active_step)
    answer = call_llm(
        worldview_profile, step_context, user_msg, passages,
        active_step=req.active_step, step_llm_guidance=step_llm_guidance,
        chat_history=history, language=user.get("language", "en"),
    )
    # Output guard: strip any handed-over deliverable blocks the model slipped in.
    answer = _strip_handed_answers(answer)

    history.append(ChatTurn(role="assistant", content=answer, step=req.active_step))
    _persist_session(sess)
    return ChatHistoryResp(session_id=req.session_id, history=history)


# ---------------- Optional streaming endpoint ----------------
@app.post("/chat/send_stream")
def chat_send_stream(req: ChatSendReq = Body(...), user: dict = Depends(get_current_user)):
    """
    Streaming variant of /chat/send — streams the LLM's answer chunk-by-chunk.
    """
    sess = _require_session(req.session_id)
    history = _get_chat(sess)

    user_msg = (req.message or "").strip()
    if not user_msg:
        raise HTTPException(status_code=400, detail="Empty message")

    # store user turn
    history.append(ChatTurn(role="user", content=user_msg, step=req.active_step))

    # Teacher-controlled mode: AI assistant may be turned off for this student's class.
    if not _student_ai_enabled(user):
        def _ai_off_stream():
            try:
                yield _canned(user, _AI_OFF_MESSAGE, _AI_OFF_MESSAGE_ES, _AI_OFF_MESSAGE_ZH)
            finally:
                history.append(ChatTurn(role="assistant", content=_canned(user, _AI_OFF_MESSAGE, _AI_OFF_MESSAGE_ES, _AI_OFF_MESSAGE_ZH), step=req.active_step))
                _persist_session(sess)
        return StreamingResponse(_ai_off_stream(), media_type="text/plain")

    # Source gate (deterministic): a citation/source request is never a safety concern,
    # so handle it before moderation to return the helpful coaching message.
    if _asks_for_sources(user_msg):
        redirect_text = _source_redirect_message(req.active_step)

        def _src_stream():
            try:
                for chunk in re.split(r"(?<=\n\n)", redirect_text):
                    if chunk:
                        yield chunk
            finally:
                history.append(ChatTurn(role="assistant", content=redirect_text, step=req.active_step))
                _persist_session(sess)

        return StreamingResponse(_src_stream(), media_type="text/plain")

    # Safety gate: refuse harmful/unethical requests (Llama Guard 3, local).
    is_safe, _cats = _moderate_input(user_msg)
    if not is_safe:
        def _refusal_stream():
            try:
                yield _canned(user, _SAFETY_REFUSAL, _SAFETY_REFUSAL_ES, _SAFETY_REFUSAL_ZH)
            finally:
                history.append(ChatTurn(role="assistant", content=_canned(user, _SAFETY_REFUSAL, _SAFETY_REFUSAL_ES, _SAFETY_REFUSAL_ZH), step=req.active_step))
                _persist_session(sess)
        return StreamingResponse(_refusal_stream(), media_type="text/plain")

    # Academic-integrity gate: block "do/rewrite it for me" requests.
    if _asks_ai_to_author(user_msg, req.active_step):
        redirect_text = _coach_redirect_message(req.active_step)

        def _redirect_stream():
            try:
                # Emit in sentence-sized chunks so it feels like a normal response
                for chunk in re.split(r"(?<=\n\n)", redirect_text):
                    if chunk:
                        yield chunk
            finally:
                history.append(ChatTurn(role="assistant", content=redirect_text, step=req.active_step))
                _persist_session(sess)

        return StreamingResponse(_redirect_stream(), media_type="text/plain")

    # Stream LLM answer
    worldview_profile = _render_worldview_profile(sess)
    step_context = _render_step_context(sess)
    passages = _retrieve(user_msg, k=5)
    step_llm_guidance = _get_step_llm_guidance(sess, req.active_step)
    payload = build_ollama_payload(
        worldview_profile, step_context, user_msg, passages,
        stream=True, active_step=req.active_step, step_llm_guidance=step_llm_guidance,
        chat_history=history, language=user.get("language", "en"),
    )

    # Capture session_id for persistence inside the generator
    session_id = sess.id

    def _stream_vllm():
        """Stream from vLLM (OpenAI SSE format)."""
        headers = {"Content-Type": "application/json"}
        if VLLM_API_KEY:
            headers["Authorization"] = f"Bearer {VLLM_API_KEY}"
        vllm_payload = {
            "model": VLLM_MODEL,
            "messages": payload["messages"],
            "temperature": LLM_TEMP,
            "max_tokens": 2048,
            "stream": True,
        }
        with requests.post(VLLM_URL, json=vllm_payload, headers=headers,
                           stream=True, timeout=300) as resp:
            resp.raise_for_status()
            for line in resp.iter_lines(decode_unicode=True):
                if not line:
                    continue
                if line.startswith("data: "):
                    line = line[6:]
                if line.strip() == "[DONE]":
                    break
                try:
                    data = json.loads(line)
                except Exception:
                    continue
                delta = data.get("choices", [{}])[0].get("delta", {}).get("content", "")
                if delta:
                    yield delta

    def _stream_ollama():
        """Stream from Ollama (native format)."""
        with requests.post(OLLAMA_URL, json=payload, stream=True, timeout=300) as resp:
            resp.raise_for_status()
            for line in resp.iter_lines(decode_unicode=True):
                if not line:
                    continue
                try:
                    data = json.loads(line)
                except Exception:
                    continue
                delta = data.get("message", {}).get("content", "")
                if delta:
                    yield delta

    def _raw_deltas():
        """Raw model token stream: try primary backend, fall back to the other."""
        stream_fn = _stream_vllm if LLM_BACKEND == "vllm" else _stream_ollama
        try:
            for delta in stream_fn():
                yield delta
        except GeneratorExit:
            raise
        except Exception as primary_err:
            logger.warning("Primary stream (%s) failed: %s — trying fallback", LLM_BACKEND, primary_err)
            fallback_fn = _stream_ollama if LLM_BACKEND == "vllm" else _stream_vllm
            for delta in fallback_fn():
                yield delta

    def event_stream():
        assistant_text_parts: List[str] = []
        try:
            # Output guard: sanitized, line-buffered stream (drops handed-over answers).
            for piece in _sanitize_stream(_raw_deltas()):
                assistant_text_parts.append(piece)
                yield piece
        except GeneratorExit:
            logger.info("Client disconnected during stream for session %s", session_id)
            return  # finally block still runs
        except Exception as e:
            logger.exception("LLM stream failed (both backends): %s", e)
            yield "\n[Error streaming from model]\n"
        finally:
            # Always persist the sanitized text, even if client disconnected mid-stream
            full_text = "".join(assistant_text_parts).strip()
            if full_text:
                history.append(ChatTurn(role="assistant", content=full_text, step=req.active_step))
            _persist_session(sess)

    return StreamingResponse(event_stream(), media_type="text/plain")


# ---------------- RAG utilities ----------------
@app.get("/rag/status")
def rag_status():
    return {
        "RAG_AVAILABLE": RAG_AVAILABLE,
        "docs_dir": str(DOCS_DIR),
        "index_path_exists": INDEX_PATH.exists(),
        "meta_path_exists": META_PATH.exists(),
        "num_chunks": len(_chunks),
    }


@app.post("/rag/reindex")
def rag_reindex():
    try:
        if INDEX_PATH.exists():
            INDEX_PATH.unlink()
        if META_PATH.exists():
            META_PATH.unlink()
    except Exception as e:
        logger.exception("Failed clearing index files: %s", e)
    _build_index()
    return {"ok": True, "num_chunks": len(_chunks)}


# ---------------- PDF Export ----------------
@app.get("/session/{session_id}/export/pdf")
def export_research_design_pdf(
    session_id: str,
    current_user: dict = Depends(get_current_user)
):
    """
    Generate and download a PDF of the research design based on all step responses.
    Automatically selects the appropriate template (quantitative, qualitative, or mixed-methods)
    based on the student's chosen methodology.
    """
    from datetime import datetime

    # Load session
    sess = _require_session(session_id)
    steps_data = sess.step_notes

    # Look up the student who owns this session (so teachers downloading get correct name)
    raw_doc = find_session(session_id)
    session_owner_id = raw_doc.get("user_id") if raw_doc else None
    if session_owner_id:
        session_owner = find_user_by_id(session_owner_id)
    else:
        session_owner = None
    # Use session owner info if available, otherwise fall back to current_user
    pdf_user = session_owner or current_user

    # Determine research path/methodology
    resolved_path = sess.resolved_path or "qualitative"  # Default to qualitative
    chosen_methodology = sess.chosen_methodology

    # Helper function to safely get field from step data
    def get_field(step_num: int, field_name: str, default: str = "") -> str:
        step_key = str(step_num)
        if step_key not in steps_data or not isinstance(steps_data[step_key], dict):
            return default
        value = steps_data[step_key].get(field_name, default)
        return str(value) if value else default

    # Extract Step 1: Worldview (+ the student's own justification, if provided)
    step1_data = steps_data.get("1", {})
    worldview = step1_data.get("worldview") or step1_data.get("worldview_id") or step1_data.get("worldview_label") or "Not specified"
    worldview_justification = (step1_data.get("worldview_justification") or "").strip()

    # Extract Step 2: Topic & Goals (detailed)
    step2_data = steps_data.get("2", {})
    step2_topic = get_field(2, "topic", "Not yet completed")
    step2_personal_goals = get_field(2, "personal_goals") or get_field(2, "personalGoals") or "Not yet completed"
    step2_practical_goals = get_field(2, "practical_goals") or get_field(2, "practicalGoals") or "Not yet completed"
    step2_intellectual_goals = get_field(2, "intellectual_goals") or get_field(2, "intellectualGoals") or "Not yet completed"

    # Extract Step 3: Conceptual Framework (detailed)
    step3_topical = get_field(3, "topicalResearch") or get_field(3, "topical_research") or "Not yet completed"
    step3_gaps = get_field(3, "gaps") or get_field(3, "gaps_identified") or "Not yet completed"
    step3_theoretical = get_field(3, "theoreticalFrameworks") or get_field(3, "theoretical_frameworks") or "Not yet completed"
    step3_problem = get_field(3, "problem_statement") or get_field(3, "problemStatement") or "Not yet completed"

    # Extract Step 4: actual research design/methodology selection (+ notes),
    # with option ids resolved to human labels.
    step4_text = _format_config_step(sess, 4, steps_data.get("4", {})) or "Not yet completed"

    # Base template data (common to all templates)
    template_data = {
        "name": pdf_user.get("username") or pdf_user.get("name", "Student"),
        "email": pdf_user.get("email") or pdf_user.get("username") or "",
        "date": datetime.now().strftime("%B %d, %Y"),
        "step1": worldview.title(),
        "step1_justification": worldview_justification,
        "step2_topic": step2_topic,
        "step2_personal_goals": step2_personal_goals,
        "step2_practical_goals": step2_practical_goals,
        "step2_intellectual_goals": step2_intellectual_goals,
        "step3_topical": step3_topical,
        "step3_gaps": step3_gaps,
        "step3_theoretical": step3_theoretical,
        "step3_problem": step3_problem,
        "step4": step4_text,
    }

    # Build Steps 5-9 straight from each step's real field keys (research aim,
    # hypothesis, collection method, participants, sampling, analysis method,
    # trustworthiness, ethics, ...), resolving option ids to labels.
    def step_text(n: int) -> str:
        return _format_config_step(sess, n, steps_data.get(str(n), {})) or "Not yet completed"

    # Mixed-methods sessions only ever fill ONE sub-methodology's fields for
    # Steps 5-9 (the pragmatist's chosen methodology), so render that path's
    # template rather than the split mixed template, which can never be populated.
    effective_path = resolved_path
    if resolved_path == "mixed" and chosen_methodology in ("quantitative", "qualitative"):
        effective_path = chosen_methodology

    if effective_path == "quantitative":
        template_name = "research_design_quantitative.html"
    elif effective_path == "qualitative":
        template_name = "research_design_qualitative.html"
    else:
        # Mixed session with no methodology chosen yet — fall back to qualitative
        # layout (Steps 5-9 will simply read "Not yet completed").
        template_name = "research_design_qualitative.html"

    # Localized template for the session owner's language (falls back to English)
    pdf_lang = (pdf_user or {}).get("language")
    if pdf_lang and pdf_lang != "en":
        loc_name = template_name.replace(".html", f".{pdf_lang}.html")
        if (TEMPLATE_DIR / loc_name).exists():
            template_name = loc_name

    for n in (5, 6, 7, 8, 9):
        template_data[f"step{n}"] = step_text(n)

    # Localize the placeholder values on a localized template
    MONTHS_ES = ["enero", "febrero", "marzo", "abril", "mayo", "junio",
                 "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre"]
    PDF_L10N = {
        "es": {
            "not_completed": "Aún no completado",
            "not_specified": "No especificado",
            "worldviews": {
                "positivist": "Positivista", "post_positivist": "Pospositivista",
                "constructivist": "Constructivista", "transformative": "Transformativa",
                "pragmatist": "Pragmatista",
            },
            "date": lambda now: f"{now.day} de {MONTHS_ES[now.month - 1]} de {now.year}",
        },
        "zh": {
            "not_completed": "尚未完成",
            "not_specified": "未指定",
            "worldviews": {
                "positivist": "实证主义", "post_positivist": "后实证主义",
                "constructivist": "建构主义", "transformative": "变革性",
                "pragmatist": "实用主义",
            },
            "date": lambda now: f"{now.year}年{now.month}月{now.day}日",
        },
    }
    tpl_lang = next((lg for lg in PDF_L10N if template_name.endswith(f".{lg}.html")), None)
    if tpl_lang:
        l10n = PDF_L10N[tpl_lang]
        for k, v in list(template_data.items()):
            if v == "Not yet completed":
                template_data[k] = l10n["not_completed"]
            elif v == "Not specified":
                template_data[k] = l10n["not_specified"]
        wv_raw = str(template_data.get("step1", "")).strip()
        wv_key = wv_raw.lower().replace(" ", "_").replace("-", "_")
        if wv_key in l10n["worldviews"]:
            template_data["step1"] = l10n["worldviews"][wv_key]
        template_data["date"] = l10n["date"](datetime.now())

    # Load and render the appropriate HTML template
    template_path = TEMPLATE_DIR / template_name
    if not template_path.exists():
        raise HTTPException(status_code=500, detail=f"PDF template not found: {template_name}")

    with open(template_path, "r", encoding="utf-8") as f:
        template_html = f.read()

    # Render with Jinja2
    template = Template(template_html)
    rendered_html = template.render(**template_data)

    # Generate PDF with WeasyPrint
    pdf_bytes = HTML(string=rendered_html, base_url=str(ROOT)).write_pdf()

    # Return PDF as download with methodology-specific filename
    methodology_name = resolved_path.title() if resolved_path else "Research"
    filename = f"{methodology_name}_Research_Design_{current_user.get('name', 'Student').replace(' ', '_')}.pdf"

    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"'
        }
    )


# ---------------- Step-data extraction helpers ----------------

def _effective_step_config(sess: "SessionData", step_num: int) -> tuple:
    """Return (effective_path, step_config) for a step, resolving mixed-methods
    inheritance and Step-4 methodology overrides to the sub-path the student
    actually used. Steps 5-9 of a mixed session inherit the chosen methodology."""
    all_paths = load_paths_config().get("paths", {})
    resolved = sess.resolved_path or "qualitative"

    effective = resolved
    if resolved == "mixed" and step_num >= 5:
        effective = sess.chosen_methodology or "qualitative"
    elif resolved != "mixed" and sess.chosen_methodology and sess.chosen_methodology != resolved:
        effective = sess.chosen_methodology

    step_cfg = all_paths.get(effective, {}).get("steps", {}).get(str(step_num), {})
    return effective, step_cfg


def _resolve_option_labels(options, value) -> str:
    """Map option id(s) to human labels. `value` may be a string id or a list of ids."""
    id_to_label = {o.get("id"): o.get("label", o.get("id")) for o in (options or [])}
    if isinstance(value, list):
        return ", ".join(str(id_to_label.get(v, v)) for v in value if v not in (None, ""))
    return str(id_to_label.get(value, value)) if value not in (None, "") else ""


def _format_config_step(sess: "SessionData", step_num: int, step_data: dict,
                        include_notes: bool = True) -> str:
    """Build a readable, labeled summary of a config-driven step (4-9) straight
    from the student's real field keys, resolving option ids to human labels.
    Returns '' if the student filled nothing in for that step."""
    if not isinstance(step_data, dict) or not step_data:
        return ""
    _, cfg = _effective_step_config(sess, step_num)
    parts = []
    ftype = cfg.get("field_type")

    if ftype in ("single_select", "multi_select"):
        label = _resolve_option_labels(cfg.get("options"), step_data.get(cfg.get("field_key")))
        if label:
            parts.append(label)
    elif ftype == "fields":
        for f in cfg.get("fields", []):
            # Respect conditional fields — skip if their depends_on condition isn't met
            dep = f.get("depends_on")
            if dep and step_data.get(dep.get("field")) != dep.get("value"):
                continue
            val = step_data.get(f.get("field_key"))
            if val in (None, "", []):
                continue
            if f.get("type") in ("select", "multi_select"):
                val = _resolve_option_labels(f.get("options"), val)
            elif isinstance(val, list):
                val = ", ".join(str(v) for v in val)
            if val:
                parts.append(f"{f.get('label', f.get('field_key'))}: {val}")
    elif ftype == "methodology_decision":
        val = step_data.get(cfg.get("field_key") or "design")
        if val:
            parts.append(f"Methodology: {val}")

    if include_notes:
        notes = step_data.get("notes")
        if notes and str(notes).strip():
            parts.append(f"Additional notes: {str(notes).strip()}")

    return "\n".join(parts)


# ---------------- PPTX Conceptual Framework Export ----------------

def _structure_cf_via_llm(sess: SessionData, raw_fields: dict) -> dict:
    """
    Call the LLM to condense/structure ONLY the fields the student has actually
    filled in.  Empty fields stay empty — the LLM must NOT invent content.
    """
    import json as _json

    # Determine which fields the student actually filled in
    filled = {k: v for k, v in raw_fields.items() if v and str(v).strip()}
    if not filled:
        # Nothing to structure — return all empty
        return raw_fields

    # Build the data block — only include filled fields
    data_lines = []
    field_labels = {
        "topic": "Topic",
        "worldview": "Worldview",
        "personal_goals": "Personal Goals",
        "topical_raw": "Topical Research (raw)",
        "theoretical_raw": "Theoretical Frameworks (raw)",
        "gaps": "Gaps",
        "problem_statement": "Problem Statement",
        "research_questions": "Research Questions",
        "research_design": "Research Design",
    }
    for key, label in field_labels.items():
        val = raw_fields.get(key, "")
        if val and str(val).strip():
            data_lines.append(f"  {label}: {str(val)[:600]}")

    data_block = "\n".join(data_lines)

    prompt = (
        "You are helping create a Conceptual Framework diagram for a research methods student.\n"
        "Condense ONLY the fields shown below into short, diagram-friendly text.\n\n"
        "CRITICAL RULE: ONLY structure fields that have data below. If a field is NOT "
        "listed in the student's data, return an EMPTY STRING for it. Do NOT invent, "
        "guess, or generate content for missing fields.\n\n"
        "FORMATTING RULES:\n"
        "- 'topics': array of up to 5 SHORT titles (3-8 words each) from their topical research. "
        "If 'Topical Research' is not provided, return an empty array [].\n"
        "- 'frameworks': array of up to 5 SHORT titles (3-8 words each) from their theoretical "
        "frameworks. If 'Theoretical Frameworks' is not provided, return an empty array [].\n"
        "- 'topic': ONE concise sentence (max 15 words). Empty string if not provided.\n"
        "- 'gaps': 1-2 sentences. Empty string if not provided.\n"
        "- 'problem_statement': 1-2 sentences. Empty string if not provided.\n"
        "- 'personal_goals': 1-2 sentences. Empty string if not provided.\n"
        "- 'research_questions': Condense but PRESERVE the research question, research aim, "
        "and hypothesis if present. Keep each clearly labeled/distinguishable (do not merge "
        "an aim and a hypothesis into one). Empty string if not provided.\n"
        "- 'research_design': 1 short sentence naming the chosen design/methodology. Empty string if not provided.\n"
        "- 'worldview': ONE word. Empty string if not provided.\n\n"
        f"STUDENT'S DATA (from 'My Research Design' panel):\n{data_block}\n\n"
        "Respond with ONLY valid JSON. No markdown, no explanation:\n"
        "{\n"
        '  "topic": "",\n'
        '  "worldview": "",\n'
        '  "personal_goals": "",\n'
        '  "topics": [],\n'
        '  "frameworks": [],\n'
        '  "gaps": "",\n'
        '  "problem_statement": "",\n'
        '  "research_questions": "",\n'
        '  "research_design": ""\n'
        "}\n"
    )

    cf_messages = [{"role": "user", "content": prompt}]
    raw = None

    # Try vLLM first if configured
    if LLM_BACKEND == "vllm":
        raw = _call_vllm(cf_messages, temperature=0.3, max_tokens=2000, timeout=90)
    if raw is None:
        # Fallback to Ollama
        ollama_result = _call_ollama({
            "model": LLM_MODEL,
            "messages": cf_messages,
            "stream": False,
            "options": {"temperature": 0.3, "num_predict": 2000},
        }, timeout=90)
        raw = ollama_result

    if not raw:
        logger.warning("Both LLM backends failed for CF structuring")
        return raw_fields

    try:
        # Extract JSON from response (handle markdown code blocks)
        json_match = re.search(r'\{[\s\S]*\}', raw)
        if json_match:
            generated = _json.loads(json_match.group())
        else:
            logger.warning("LLM did not return valid JSON for CF: %s", raw[:200])
            return raw_fields

        # Ensure topics and frameworks are lists of 5
        if isinstance(generated.get("topics"), list):
            generated["topics"] = (generated["topics"] + [""] * 5)[:5]
        if isinstance(generated.get("frameworks"), list):
            generated["frameworks"] = (generated["frameworks"] + [""] * 5)[:5]

        return generated
    except Exception as e:
        logger.warning("LLM CF structuring failed: %s", e)
        return raw_fields


def _gather_cf_data(session_id: str, current_user: dict) -> dict:
    """Shared helper: gather conceptual framework data, always using LLM to structure."""
    from datetime import datetime

    sess = _require_session(session_id)
    steps_data = sess.step_notes

    step1_data = steps_data.get("1", {})
    worldview = (step1_data.get("worldview") or step1_data.get("worldview_id")
                 or step1_data.get("worldview_label") or "")

    step2_data = steps_data.get("2", {})
    topic = step2_data.get("topic", "")
    # Three separate goal fields (new) with fallback to old single "goals" field
    pg = step2_data.get("personalGoals") or step2_data.get("personal_goals") or ""
    pr = step2_data.get("practicalGoals") or step2_data.get("practical_goals") or ""
    ig = step2_data.get("intellectualGoals") or step2_data.get("intellectual_goals") or ""
    # Combine for CF — or fall back to old single field
    if pg or pr or ig:
        goal_parts = []
        if pg: goal_parts.append(f"Personal: {pg}")
        if pr: goal_parts.append(f"Practical: {pr}")
        if ig: goal_parts.append(f"Intellectual: {ig}")
        personal_goals = "; ".join(goal_parts)
    else:
        personal_goals = step2_data.get("goals", "")

    step3_data = steps_data.get("3", {})
    topical_raw = step3_data.get("topicalResearch") or step3_data.get("topical_research") or ""
    theoretical_raw = step3_data.get("theoreticalFrameworks") or step3_data.get("theoretical_frameworks") or ""
    gaps = step3_data.get("gaps") or step3_data.get("gaps_identified") or ""
    problem = step3_data.get("problem_statement") or step3_data.get("problemStatement") or ""

    # Step 4: the actual selected research design/methodology (+ any notes),
    # not just the free-text "additional questions" box.
    step4_data = steps_data.get("4", {})
    research_design = _format_config_step(sess, 4, step4_data)

    # Step 5: research question, aim, and hypothesis — using the real field keys.
    # Quantitative students enter research_aim + hypothesis; qualitative students
    # enter research_question. Preserve each so aims/hypotheses map through clearly.
    step5_data = steps_data.get("5", {})
    research_questions = _format_config_step(sess, 5, step5_data)
    if not research_questions:
        research_questions = (step5_data.get("research_question") or step5_data.get("notes") or "")

    # Attribute the export to the SESSION OWNER (the student), not whoever is
    # downloading it — a teacher/admin viewing a student's design must not have
    # their own name/email stamped on the student's conceptual framework.
    raw_doc = find_session(session_id)
    owner_id = raw_doc.get("user_id") if raw_doc else None
    owner = find_user_by_id(owner_id) if owner_id else None
    export_user = owner or current_user
    email = export_user.get("email") or export_user.get("username") or ""
    name = export_user.get("name", "Student")
    timestamp = datetime.now().strftime("%B %d, %Y")

    # Pass raw text to LLM — it will extract short titles and structure everything
    raw_fields = {
        "topic": topic,
        "worldview": worldview,
        "personal_goals": personal_goals,
        "topical_raw": topical_raw,
        "theoretical_raw": theoretical_raw,
        "gaps": gaps,
        "problem_statement": problem,
        "research_questions": research_questions,
        "research_design": research_design,
    }

    # Structure via LLM, but cache the result keyed on the raw inputs: the
    # model runs once per unique set of step notes, so reopening the editor
    # is instant unless the student actually changed their notes.
    import hashlib
    raw_hash = hashlib.md5(
        json.dumps(raw_fields, sort_keys=True).encode("utf-8")
    ).hexdigest()
    cache = (raw_doc or {}).get("cf_prefill_cache") or {}
    if cache.get("hash") == raw_hash and isinstance(cache.get("structured"), dict):
        structured = cache["structured"]
    else:
        structured = _structure_cf_via_llm(sess, raw_fields)
        update_session(session_id, {"cf_prefill_cache": {"hash": raw_hash, "structured": structured}})

    # Extract topics/frameworks — only if the student wrote topical/theoretical data
    topics = []
    if topical_raw.strip():
        topics = structured.get("topics", [])
        if not topics:
            lines = [l.strip(" -•·\t") for l in topical_raw.split("\n") if l.strip(" -•·\t")]
            topics = lines[:5]
    topics = (topics + [""] * 5)[:5]

    frameworks = []
    if theoretical_raw.strip():
        frameworks = structured.get("frameworks", [])
        if not frameworks:
            lines = [l.strip(" -•·\t") for l in theoretical_raw.split("\n") if l.strip(" -•·\t")]
            frameworks = lines[:5]
    frameworks = (frameworks + [""] * 5)[:5]

    # Only use structured values for fields the student actually filled in
    wv = (structured.get("worldview") if worldview.strip() else "") or ""
    return {
        "email": email,
        "name": name,
        "date": timestamp,
        "topic": (structured.get("topic") if topic.strip() else "") or "",
        "worldview": wv.title() if wv else "",
        "personal_goals": (structured.get("personal_goals") if personal_goals.strip() else "") or "",
        "topics": topics,
        "frameworks": frameworks,
        "gaps": (structured.get("gaps") if gaps.strip() else "") or "",
        "problem_statement": (structured.get("problem_statement") if problem.strip() else "") or "",
        # Prefer the LLM-condensed text, but fall back to the student's raw
        # labeled entries so aims/hypotheses/design are never silently dropped.
        "research_questions": ((structured.get("research_questions") or research_questions) if research_questions.strip() else "") or "",
        "research_design": ((structured.get("research_design") or research_design) if research_design.strip() else "") or "",
    }


# Scalar fields the CF editor can save; topics/frameworks are 5-item lists.
CF_FIELD_KEYS = [
    "topic", "worldview", "personal_goals", "gaps", "problem_statement",
    "research_questions", "research_design",
]


class ConceptualFrameworkSaveRequest(BaseModel):
    fields: Dict[str, Any]


def _cf_sanitize(fields: dict) -> dict:
    """Whitelist + clamp the editor payload before persisting."""
    clean = {}
    for k in CF_FIELD_KEYS:
        if k in fields:
            clean[k] = str(fields.get(k) or "")[:2000]
    for k in ("topics", "frameworks"):
        if k in fields:
            vals = fields.get(k) or []
            if isinstance(vals, list):
                clean[k] = [str(v or "")[:500] for v in vals[:5]]
    return clean


def _cf_effective_data(session_id: str, current_user: dict) -> dict:
    """Saved editor fields win; otherwise prefill via the LLM gatherer.
    When a save exists we skip the LLM entirely (fast reopen, and the
    student's edits are authoritative). Identity/date always fresh."""
    raw_doc = find_session(session_id)
    stored = (raw_doc or {}).get("conceptual_framework_fields") or {}
    if stored:
        base = _gather_cf_identity(session_id, current_user)
        base["topics"] = ((stored.get("topics") or []) + [""] * 5)[:5]
        base["frameworks"] = ((stored.get("frameworks") or []) + [""] * 5)[:5]
        for k in CF_FIELD_KEYS:
            base[k] = str(stored.get(k) or "")
        base["has_saved"] = True
        return base
    data = _gather_cf_data(session_id, current_user)
    data["has_saved"] = False
    return data


def _gather_cf_identity(session_id: str, current_user: dict) -> dict:
    """Just the owner identity + date (no LLM), same attribution rules as
    _gather_cf_data."""
    from datetime import datetime
    _require_session(session_id)
    raw_doc = find_session(session_id)
    owner_id = raw_doc.get("user_id") if raw_doc else None
    owner = find_user_by_id(owner_id) if owner_id else None
    export_user = owner or current_user
    return {
        "email": export_user.get("email") or export_user.get("username") or "",
        "name": export_user.get("name", "Student"),
        "date": datetime.now().strftime("%B %d, %Y"),
    }


@app.get("/session/{session_id}/export/conceptual-framework/data")
def get_conceptual_framework_data(
    session_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Return conceptual framework data as JSON for the web editor
    (saved editor fields win over the LLM-structured prefill)."""
    return _cf_effective_data(session_id, current_user)


@app.put("/session/{session_id}/conceptual-framework/data")
def save_conceptual_framework_data(
    session_id: str,
    payload: ConceptualFrameworkSaveRequest,
    current_user: dict = Depends(get_current_user)
):
    """Persist the student's conceptual framework edits (from the editor)."""
    _require_session(session_id)
    update_session(session_id, {"conceptual_framework_fields": _cf_sanitize(payload.fields)})
    return {"ok": True}


@app.get("/session/{session_id}/export/conceptual-framework")
def export_conceptual_framework(
    session_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Generate a Conceptual Framework PPTX using the shared data helper."""
    from pptx import Presentation as PptxPresentation
    import io

    d = _cf_effective_data(session_id, current_user)
    topics = (d.get("topics", []) + [""] * 5)[:5]
    frameworks = (d.get("frameworks", []) + [""] * 5)[:5]

    replacements = {
        "<<email>>": d["email"],
        "<<Timestamp>>": d["date"],
        "<<Define your research topic>>": d["topic"] or "Not yet defined",
        "<<Name>>": d["name"],
        "<<Gap/s found in the review of your topical research>>": d["gaps"] or "Not yet identified",
        "<<Define your Personal Interests and Goals>>": d["personal_goals"] or "Not yet defined",
        "<<Describe your positionality and worldview >>": d["worldview"] or "Not specified",
        "<<Define your Problem Statement>>": d["problem_statement"] or "Not yet defined",
        "<<Define your Research Question/s>>": d["research_questions"] or "Not yet defined",
        "<<Define your Research Design>>": d["research_design"] or "Not yet defined",
    }
    for i in range(5):
        replacements[f"<<Topic {i+1}>>"] = topics[i] if topics[i] else ""
        replacements[f"<<Theoretical Framework {i+1}>>"] = frameworks[i] if frameworks[i] else ""

    template_path = TEMPLATE_DIR / "conceptual_framework.pptx"
    if not template_path.exists():
        raise HTTPException(status_code=500, detail="Conceptual framework template not found")

    prs = PptxPresentation(str(template_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_pptx_text(shape, replacements)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"Conceptual_Framework_{d['name'].replace(' ', '_')}.pptx"
    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


# ---------------- PPTX Visual Design Export (qualitative) ----------------

# Which slide (0-based) of visual_design_qualitative.pptx belongs to each
# Step-4 qualitative design. Slide 0 is the generic legend slide.
# school_ethnography reuses the ethnography layout; design_based_research
# reuses the action-research layout (both are cycle-based improvement designs).
VD_SLIDE_BY_DESIGN = {
    "case_study": 1,
    "action_research": 2,
    "design_based_research": 2,
    "ethnography": 3,
    "school_ethnography": 3,
    "narrative": 4,
    "phenomenology": 5,
    "phenomenography": 6,
    "grounded_theory": 7,
}

VD_CENTRAL_LABEL = {
    "case_study": "the case being studied",
    "action_research": "the aspect to reflect upon and improve",
    "design_based_research": "the aspect to reflect upon and improve",
    "ethnography": "the culture/group under study",
    "school_ethnography": "the culture/group under study",
    "narrative": "the narrative portraits",
    "phenomenology": "the phenomenon under study",
    "phenomenography": "the phenomenon under study",
    "grounded_theory": "the phenomenon under study (substantive area of interest)",
}


def _structure_vd_via_llm(design_label: str, central_label: str, raw_fields: dict) -> dict:
    """
    Ask the LLM to condense the student's step data into the short,
    diagram-friendly snippets the visual-design slide needs. Only fields the
    student actually filled in are sent; the LLM must not invent content.
    """
    import json as _json

    filled = {k: v for k, v in raw_fields.items() if v and str(v).strip()}
    if not filled:
        return {}

    data_lines = [f"  {k}: {str(v)[:600]}" for k, v in filled.items()]
    data_block = "\n".join(data_lines)

    prompt = (
        f"You are filling in a one-slide visual diagram of a student's {design_label} "
        "research design. Condense ONLY the student data below into short, "
        "diagram-friendly text (each value at most 25 words; fragments are fine).\n\n"
        "CRITICAL RULE: only use information present in the data below. If nothing in "
        "the data answers a field, return an empty string for it. Do NOT invent content.\n\n"
        "Fields to produce:\n"
        f"- 'central_item': {central_label}, taken from their topic/question.\n"
        "- 'context': the setting/context where the study takes place.\n"
        "- 'informants': who the participants/informants are.\n"
        "- 'other_documents': documents/artifacts to be analyzed besides direct data collection.\n"
        "- 'data_gathering': the data gathering methods.\n"
        "- 'strategies': design/analysis strategies (sampling, cycles, analysis approach).\n"
        "- 'process_support': supports for rigor (trustworthiness strategies, ethics safeguards, tools).\n"
        "- 'question': the central research question or issue, lightly condensed.\n"
        "- 'topics': up to 4 short topic titles from their topical research, joined with newlines.\n"
        "- 'minicases': mini-cases or sub-units of analysis, if any.\n\n"
        f"STUDENT'S DATA:\n{data_block}\n\n"
        "Respond with ONLY valid JSON, no markdown:\n"
        "{\n"
        '  "central_item": "", "context": "", "informants": "", "other_documents": "",\n'
        '  "data_gathering": "", "strategies": "", "process_support": "", "question": "",\n'
        '  "topics": "", "minicases": ""\n'
        "}\n"
    )

    vd_messages = [{"role": "user", "content": prompt}]
    raw = None
    if LLM_BACKEND == "vllm":
        raw = _call_vllm(vd_messages, temperature=0.3, max_tokens=1200, timeout=90)
    if raw is None:
        raw = _call_ollama({
            "model": LLM_MODEL,
            "messages": vd_messages,
            "stream": False,
            "options": {"temperature": 0.3, "num_predict": 1200},
        }, timeout=90)

    if not raw:
        logger.warning("Both LLM backends failed for visual design structuring")
        return {}
    try:
        json_match = re.search(r'\{[\s\S]*\}', raw)
        if not json_match:
            logger.warning("LLM did not return valid JSON for visual design: %s", raw[:200])
            return {}
        parsed = _json.loads(json_match.group())
        return parsed if isinstance(parsed, dict) else {}
    except Exception as e:
        logger.warning("LLM visual design structuring failed: %s", e)
        return {}


VD_FIELD_KEYS = [
    # shared / qualitative
    "central_item", "context", "question", "topics", "informants",
    "data_gathering", "other_documents", "strategies", "process_support", "minicases",
    # quantitative
    "variables", "sample", "groups", "data_analysis", "study_type",
    # quantitative continuum sliders (0-100, stored as strings; empty = design default)
    "slider_variance", "slider_causality", "slider_iv_control",
    # mixed methods (second-strand keys; the qualitative strand reuses the
    # qualitative keys and the quantitative strand reuses the quantitative ones)
    "research_topic", "mm_question", "mm_data_gathering", "mm_process_support",
    "qual_tradition", "hypothesis", "qual_question", "embedded_host",
]

# Quantitative designs supported by the visual design editor/export.
# The quantitative template's fill-in slide is slide 2 (index 1).
VD_QUANT_DESIGNS = {"descriptive", "correlational", "quasi_experimental", "experimental", "cross_sectional_survey", "pre_experimental"}

# Mixed methods designs supported by the visual design editor (print-only;
# stored in step_notes["4"]["mixed_design"] for pragmatist/mixed sessions)
VD_MIXED_DESIGNS = {
    "convergent_parallel": "Convergent Parallel Mixed Methods",
    "explanatory_sequential": "Explanatory Sequential Mixed Methods",
    "exploratory_sequential": "Exploratory Sequential Mixed Methods",
    "embedded": "Embedded Mixed Methods",
}
VD_MIXED_EDITOR_READY = {"convergent_parallel", "explanatory_sequential", "exploratory_sequential", "embedded"}


def _vd_context(session_id: str, current_user: dict):
    """Shared validation + owner lookup for the visual design endpoints.
    Returns (sess, raw_doc, design_id, design_label, name, email)."""
    sess = _require_session(session_id)
    steps_data = sess.step_notes

    resolved = sess.resolved_path or "qualitative"
    effective = resolved
    if resolved == "mixed":
        effective = sess.chosen_methodology or ""
    elif sess.chosen_methodology:
        effective = sess.chosen_methodology

    design_id = (steps_data.get("4") or {}).get("design")
    mixed_design = (steps_data.get("4") or {}).get("mixed_design")

    # The design saved at Step 4 is the ground truth for which template family
    # applies: a stale chosen_methodology (e.g. the student explored the
    # Step-4 methodology override and then went back) must not route a
    # qualitative design down the quantitative path or vice versa.
    if not (resolved == "mixed" and mixed_design):
        if design_id in VD_SLIDE_BY_DESIGN:
            effective = "qualitative"
        elif design_id in VD_QUANT_DESIGNS:
            effective = "quantitative"

    if resolved == "mixed" and mixed_design:
        # Pragmatist student who chose a mixed methods design: the visual
        # design combines both strands regardless of the primary methodology.
        if mixed_design not in VD_MIXED_DESIGNS:
            raise HTTPException(status_code=400, detail="Unknown mixed methods design.")
        if mixed_design not in VD_MIXED_EDITOR_READY:
            raise HTTPException(
                status_code=400,
                detail="The visual design is not available for this mixed methods design yet."
            )
        effective = "mixed"
        design_id = mixed_design
    elif effective == "qualitative":
        if design_id not in VD_SLIDE_BY_DESIGN:
            raise HTTPException(
                status_code=400,
                detail="Complete Step 4 (choose your qualitative design) to work on the visual design."
            )
    elif effective == "quantitative":
        if not design_id:
            raise HTTPException(
                status_code=400,
                detail="Complete Step 4 (choose your quantitative design) to work on the visual design."
            )
        if design_id not in VD_QUANT_DESIGNS:
            raise HTTPException(
                status_code=400,
                detail="The visual design is not available for this quantitative design yet."
            )
    else:
        raise HTTPException(
            status_code=400,
            detail="Choose your methodology in Step 4 to work on the visual design."
        )

    # Session owner info (teachers get the student's name, not their own)
    raw_doc = find_session(session_id)
    owner = find_user_by_id(raw_doc.get("user_id")) if raw_doc and raw_doc.get("user_id") else None
    vd_user = owner or current_user
    name = vd_user.get("username") or vd_user.get("name", "Student")
    email = vd_user.get("email") or vd_user.get("username") or ""

    if effective == "mixed":
        design_label = VD_MIXED_DESIGNS[design_id]
    else:
        design_label = _resolve_option_labels(
            _effective_step_config(sess, 4)[1].get("options"), design_id) or design_id
    return sess, raw_doc, design_id, design_label, name, email, effective


def _vd_raw_fields(sess: SessionData) -> dict:
    """The student's step data relevant to the visual design, labels resolved."""
    steps_data = sess.step_notes

    def get_field(step_num: int, *keys) -> str:
        d = steps_data.get(str(step_num))
        if not isinstance(d, dict):
            return ""
        for k in keys:
            v = d.get(k)
            if v not in (None, "", []):
                return _resolve_option_labels(None, v) if isinstance(v, list) else str(v)
        return ""

    def get_labeled(step_num: int, field_key: str) -> str:
        d = steps_data.get(str(step_num))
        if not isinstance(d, dict):
            return ""
        val = d.get(field_key)
        if val in (None, "", []):
            return ""
        _, cfg = _effective_step_config(sess, step_num)
        if cfg.get("field_type") in ("single_select", "multi_select") and cfg.get("field_key") == field_key:
            return _resolve_option_labels(cfg.get("options"), val)
        for f in cfg.get("fields", []):
            if f.get("field_key") == field_key and f.get("type") in ("select", "multi_select"):
                return _resolve_option_labels(f.get("options"), val)
        return ", ".join(str(v) for v in val) if isinstance(val, list) else str(val)

    return {
        "Topic": get_field(2, "topic"),
        "Topical Research": get_field(3, "topicalResearch", "topical_research"),
        "Research Question / Central Issue": get_labeled(5, "research_question"),
        "Research Aim": get_labeled(5, "research_aim"),
        "Hypothesis": get_labeled(5, "hypothesis"),
        "Collection Method": get_labeled(6, "collection_method"),
        "Participants": get_labeled(6, "participants"),
        "Sampling Method": get_labeled(6, "sampling_method"),
        "Analysis Method": get_labeled(7, "analysis_method"),
        "Analysis Plan": get_labeled(7, "analysis_notes"),
        "Trustworthiness Strategies": get_labeled(8, "trustworthiness_methods"),
        "Ethics Plan": get_labeled(9, "ethics_approach") or get_labeled(9, "ethics_notes"),
    }


def _vd_prefill(raw_fields: dict) -> dict:
    """Direct (no-LLM) prefill of the visual design fields from step data."""
    topics_fallback = "\n".join(
        l.strip(" -•·\t") for l in raw_fields["Topical Research"].split("\n") if l.strip(" -•·\t")
    )[:300]
    sample_parts = [p for p in (raw_fields["Participants"].strip(), raw_fields["Sampling Method"].strip()) if p]
    return {
        "central_item": raw_fields["Topic"].strip(),
        "context": "",
        "question": (raw_fields["Research Question / Central Issue"] or raw_fields["Research Aim"]).strip(),
        "topics": topics_fallback,
        "informants": raw_fields["Participants"].strip(),
        "data_gathering": raw_fields["Collection Method"].strip(),
        "other_documents": "",
        "strategies": (raw_fields["Analysis Method"] or raw_fields["Sampling Method"]).strip(),
        "process_support": raw_fields["Trustworthiness Strategies"].strip(),
        "minicases": "",
        # quantitative
        "variables": raw_fields["Hypothesis"].strip(),
        "sample": ". ".join(sample_parts),
        "groups": "",
        "data_analysis": raw_fields["Analysis Method"].strip(),
        "study_type": "",
        "slider_variance": "",
        "slider_causality": "",
        "slider_iv_control": "",
        "research_topic": raw_fields["Topic"].strip(),
        "mm_question": raw_fields["Research Aim"].strip(),
        "mm_data_gathering": "",
        "mm_process_support": "",
        "qual_tradition": "",
        "hypothesis": raw_fields["Hypothesis"].strip(),
        "qual_question": raw_fields["Research Question / Central Issue"].strip(),
        "embedded_host": "",
    }


class VisualDesignSaveRequest(BaseModel):
    fields: Dict[str, str]


@app.get("/session/{session_id}/visual-design/data")
def get_visual_design_data(
    session_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Data for the visual design editor page: the student's saved fields,
    prefilled from their step data where nothing was saved yet."""
    sess, raw_doc, design_id, design_label, name, email, effective = _vd_context(session_id, current_user)

    stored = (raw_doc or {}).get("visual_design_fields") or {}
    prefill = _vd_prefill(_vd_raw_fields(sess))
    fields = {}
    for key in VD_FIELD_KEYS:
        saved_val = stored.get(key)
        fields[key] = str(saved_val).strip() if saved_val not in (None, "") else prefill.get(key, "")

    return {
        "design": design_id,
        "design_label": design_label,
        "path": effective,
        "primary": ((sess.step_notes.get("4") or {}).get("embedded_host")
                    or sess.chosen_methodology or "qualitative"),
        "name": name,
        "email": email,
        "fields": fields,
        "has_saved": bool(stored),
    }


@app.put("/session/{session_id}/visual-design/data")
def save_visual_design_data(
    session_id: str,
    payload: VisualDesignSaveRequest,
    current_user: dict = Depends(get_current_user)
):
    """Persist the student's visual design fields (from the editor page)."""
    _vd_context(session_id, current_user)
    fields = {
        k: str(payload.fields.get(k, ""))[:2000]
        for k in VD_FIELD_KEYS if k in payload.fields
    }
    update_session(session_id, {"visual_design_fields": fields})
    return {"ok": True}




# ============================================================
# Admin endpoints
# ============================================================


@app.get("/admin/stats")
def admin_stats(admin: dict = Depends(require_admin)):
    """Aggregate stats for admin dashboard."""
    role_counts = get_user_counts_by_role()
    total_users = sum(role_counts.values())
    return {
        "total_users": total_users,
        "role_counts": role_counts,
        "total_sessions": get_total_sessions_count(),
        "total_classes": get_total_classes_count(),
        "active_users_7d": get_active_users_last_n_days(7),
        "active_users_30d": get_active_users_last_n_days(30),
    }


@app.get("/admin/signups")
def admin_signups(days: int = Query(30), admin: dict = Depends(require_admin)):
    """Daily signup counts for chart."""
    return {"signups": get_signups_over_time(days)}


@app.get("/admin/step-completion")
def admin_step_completion(admin: dict = Depends(require_admin)):
    """Step completion across all students."""
    return {"steps": get_step_completion_across_all()}


@app.get("/admin/login-activity")
def admin_login_activity(
    limit: int = Query(100),
    skip: int = Query(0),
    search: str = Query(""),
    status: str = Query(""),  # "ok" | "fail" | ""
    admin: dict = Depends(require_admin),
):
    """Recent login history log with search and OK/FAIL filter."""
    success = {"ok": True, "fail": False}.get(status)
    logins, total = get_recent_logins(limit, skip, search or None, success)
    return {"logins": logins, "total": total}


@app.get("/admin/login-map")
def admin_login_map(days: int = Query(0), admin: dict = Depends(require_admin)):
    """Geo-located login data for world map (days=0 means all time)."""
    return {"locations": get_login_locations(days)}


@app.get("/admin/login-timeseries")
def admin_login_timeseries(days: int = Query(0), admin: dict = Depends(require_admin)):
    """Logins per day, success/failed split (days=0 means all time)."""
    from database import get_login_timeseries
    return {"series": get_login_timeseries(days)}


@app.get("/admin/geo/countries")
def admin_geo_countries(days: int = Query(0), admin: dict = Depends(require_admin)):
    """User/login counts aggregated by country (days=0 means all time)."""
    return {"countries": get_login_stats_by_country(days)}


@app.get("/admin/geo/regions")
def admin_geo_regions(days: int = Query(0), admin: dict = Depends(require_admin)):
    """User/login counts aggregated by city/region (days=0 means all time)."""
    return {"regions": get_login_stats_by_region(days)}


@app.get("/admin/users")
def admin_list_users(
    skip: int = Query(0),
    limit: int = Query(50),
    role: Optional[str] = Query(None),
    search: Optional[str] = Query(None),
    sort_by: str = Query("created_at"),
    sort_dir: str = Query("desc"),
    admin: dict = Depends(require_admin),
):
    """Paginated user list with optional role filter, search, and sorting."""
    users, total = get_all_users(skip, limit, role, search, sort_by, sort_dir)
    for u in users:
        u["_id"] = str(u["_id"])
    return {"users": users, "total": total}


@app.get("/admin/users/{user_id}")
def admin_get_user(user_id: str, admin: dict = Depends(require_admin)):
    """Single user detail + recent logins."""
    user = find_user_by_id(user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    user["_id"] = str(user["_id"])
    user.pop("password_hash", None)
    logins = get_logins_for_user(user_id, limit=20)
    return {"user": user, "logins": logins}


class AdminCreateUserReq(BaseModel):
    email: str
    password: str
    name: str
    role: str
    education_level: str = "high_school"


@app.post("/admin/users")
def admin_create_user(req: AdminCreateUserReq, admin: dict = Depends(require_admin)):
    """Create a new user (any role including admin)."""
    if req.role not in ("student", "teacher", "admin"):
        raise HTTPException(status_code=400, detail="Invalid role")
    if find_user_by_email(req.email):
        raise HTTPException(status_code=409, detail="Email already registered")
    pw_hash = hash_password(req.password)
    user_id = create_user(req.email, pw_hash, req.role, req.name, req.education_level)
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "create_user", user_id, req.email, {"role": req.role}
    )
    return {"ok": True, "user_id": user_id}


class AdminUpdateUserReq(BaseModel):
    name: Optional[str] = None
    role: Optional[str] = None
    education_level: Optional[str] = None
    is_active: Optional[bool] = None


@app.patch("/admin/users/{user_id}")
def admin_update_user(user_id: str, req: AdminUpdateUserReq, admin: dict = Depends(require_admin)):
    """Update user fields."""
    user = find_user_by_id(user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    fields = {}
    details = {}
    if req.name is not None:
        fields["name"] = req.name
        details["name"] = req.name
    if req.role is not None:
        if req.role not in ("student", "teacher", "admin", "classroom_student"):
            raise HTTPException(status_code=400, detail="Invalid role")
        details["old_role"] = user.get("role")
        details["new_role"] = req.role
        fields["role"] = req.role
    if req.education_level is not None:
        fields["education_level"] = req.education_level
    if req.is_active is not None:
        fields["is_active"] = req.is_active
        details["is_active"] = req.is_active
    if not fields:
        raise HTTPException(status_code=400, detail="No fields to update")
    update_user_fields(user_id, fields)
    action = "update_user"
    if req.is_active is False:
        action = "deactivate"
    elif req.is_active is True:
        action = "reactivate"
    if req.role and req.role != user.get("role"):
        action = "change_role"
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        action, user_id, user.get("email") or user.get("username", ""), details
    )
    return {"ok": True}


class AdminResetPasswordReq(BaseModel):
    new_password: str


@app.post("/admin/users/{user_id}/reset-password")
def admin_reset_password(user_id: str, req: AdminResetPasswordReq, admin: dict = Depends(require_admin)):
    """Admin force-resets a user's password."""
    user = find_user_by_id(user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    if len(req.new_password) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters")
    new_hash = hash_password(req.new_password)
    email = user.get("email")
    if email:
        update_user_password(email, new_hash)
    else:
        update_user_fields(user_id, {"password_hash": new_hash})
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "reset_password", user_id,
        user.get("email") or user.get("username", ""), {}
    )
    return {"ok": True}


@app.delete("/admin/users/{user_id}")
def admin_delete_user(user_id: str, admin: dict = Depends(require_admin)):
    """Hard-delete a user."""
    user = find_user_by_id(user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    if str(admin["_id"]) == user_id:
        raise HTTPException(status_code=400, detail="Cannot delete your own account")
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "delete_user", user_id,
        user.get("email") or user.get("username", ""),
        {"role": user.get("role"), "name": user.get("name")}
    )
    delete_user_by_id(user_id)
    return {"ok": True}


@app.get("/admin/audit-log")
def admin_audit_log(
    limit: int = Query(100),
    skip: int = Query(0),
    action: str = Query(""),
    admin: dict = Depends(require_admin),
):
    from database import get_audit_action_types
    log, total = get_admin_audit_log(limit, skip, action or None)
    return {"log": log, "total": total, "actions": get_audit_action_types()}


# ── Glossary ────────────────────────────────────────────

@app.get("/glossary")
def public_glossary(lang: Optional[str] = Query(None)):
    """All glossary terms — read by the student Dictionary tab. Public.
    Spanish terms are served from stored auto-translations, falling back to
    English per term while a translation is pending."""
    lang = lang if lang in SUPPORTED_LANGUAGES else "en"
    return {"terms": get_all_glossary_terms(lang), "lang": lang}


GLOSSARY_LANG_NAMES = {"es": "Spanish", "zh": "Simplified Chinese"}


def _translate_glossary_term(term: str, definition: str, lang: str = "es") -> Optional[dict]:
    """One-shot LLM translation of a glossary entry into the target language.
    Returns {'term', 'def'} or None if both backends fail."""
    lang_name = GLOSSARY_LANG_NAMES.get(lang, "Spanish")
    prompt = (
        f"You translate research-methods glossary entries from English to {lang_name} "
        f"for high-school and university students. Use the standard {lang_name} term "
        "used in research-methodology courses (not a literal word-for-word "
        "rendering). Keep the definition's plain, student-friendly tone. If the "
        "English term includes a parenthetical, keep an equivalent parenthetical.\n\n"
        f"TERM: {term}\n"
        f"DEFINITION: {definition}\n\n"
        "Respond with ONLY valid JSON, no markdown, no explanation:\n"
        '{"term": "", "def": ""}'
    )
    messages = [{"role": "user", "content": prompt}]
    raw = None
    if LLM_BACKEND == "vllm":
        raw = _call_vllm(messages, temperature=0.2, max_tokens=500, timeout=60)
    if raw is None:
        raw = _call_ollama({
            "model": LLM_MODEL,
            "messages": messages,
            "stream": False,
            "options": {"temperature": 0.2, "num_predict": 500},
        }, timeout=60)
    if not raw:
        return None
    try:
        match = re.search(r'\{[\s\S]*\}', raw)
        data = _json.loads(match.group()) if match else None
        if data and (data.get("term") or "").strip() and (data.get("def") or "").strip():
            return {"term": data["term"].strip(), "def": data["def"].strip()}
    except Exception as e:
        logger.warning("Glossary translation parse failed for %r: %s", term, e)
    return None


def _translate_glossary_ids(term_ids: List[str], langs: tuple = ("es", "zh")):
    """Background worker: translate the given glossary terms into every overlay
    language that still lacks one. Skips terms already translated (or deleted)
    by the time it runs."""
    done = failed = 0
    for tid in term_ids:
        doc = get_glossary_term_by_id(tid)
        if not doc:
            continue
        for lg in langs:
            if doc.get(f"term_{lg}") and doc.get(f"def_{lg}"):
                continue
            result = _translate_glossary_term(doc.get("term", ""), doc.get("def", ""), lg)
            if result and set_glossary_translation(tid, result["term"], result["def"], lg):
                done += 1
            else:
                failed += 1
    if done or failed:
        logger.info("[glossary] Auto-translated %d entry/entries (%d failed)", done, failed)


class GlossaryTermReq(BaseModel):
    term: Optional[str] = None
    definition: Optional[str] = None
    steps: Optional[List[int]] = None


@app.get("/admin/glossary")
def admin_list_glossary(admin: dict = Depends(require_admin)):
    terms = get_all_glossary_terms()
    return {
        "terms": terms,
        "total": count_glossary_terms(),
        "missing_es": sum(1 for t in terms if not t.get("has_es")),
        "missing_zh": sum(1 for t in terms if not t.get("has_zh")),
    }


@app.post("/admin/glossary")
def admin_create_glossary(req: GlossaryTermReq, background_tasks: BackgroundTasks,
                          admin: dict = Depends(require_admin)):
    if not (req.term or "").strip() or not (req.definition or "").strip():
        raise HTTPException(status_code=400, detail="Term and definition are required")
    steps = [int(s) for s in (req.steps or []) if 1 <= int(s) <= 9]
    term_id = create_glossary_term(req.term, req.definition, steps)
    background_tasks.add_task(_translate_glossary_ids, [term_id])
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "create_glossary_term", term_id, "", {"term": req.term.strip()}
    )
    return {"ok": True, "id": term_id}


@app.patch("/admin/glossary/{term_id}")
def admin_update_glossary(term_id: str, req: GlossaryTermReq, background_tasks: BackgroundTasks,
                          admin: dict = Depends(require_admin)):
    fields: dict = {}
    if req.term is not None:
        fields["term"] = req.term
    if req.definition is not None:
        fields["def"] = req.definition
    if req.steps is not None:
        fields["steps"] = [int(s) for s in req.steps if 1 <= int(s) <= 9]
    if not fields:
        raise HTTPException(status_code=400, detail="No fields to update")
    ok = update_glossary_term(term_id, fields)
    if not ok:
        raise HTTPException(status_code=404, detail="Term not found")
    # English text changed → stored Spanish was cleared; refresh it
    if "term" in fields or "def" in fields:
        background_tasks.add_task(_translate_glossary_ids, [term_id])
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "update_glossary_term", term_id, "", {"term": fields.get("term")}
    )
    return {"ok": True}


@app.post("/admin/glossary/translate-missing")
def admin_glossary_translate_missing(background_tasks: BackgroundTasks,
                                     admin: dict = Depends(require_admin)):
    """Queue translations for every glossary term missing any overlay language."""
    ids = sorted(set(get_glossary_ids_missing("es")) | set(get_glossary_ids_missing("zh")))
    if ids:
        background_tasks.add_task(_translate_glossary_ids, ids)
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "translate_glossary", "", "", {"queued": len(ids)}
    )
    return {"ok": True, "queued": len(ids)}


@app.delete("/admin/glossary/{term_id}")
def admin_delete_glossary(term_id: str, admin: dict = Depends(require_admin)):
    ok = delete_glossary_term(term_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Term not found")
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "delete_glossary_term", term_id, "", {}
    )
    return {"ok": True}


# ── Knowledge base / RAG resources ──────────────────────

def _safe_resource_name(filename: str) -> str:
    """Sanitize an uploaded filename to a bare, safe basename."""
    name = os.path.basename(filename or "").strip()
    if not name or name in (".", "..") or "/" in name or "\\" in name:
        raise HTTPException(status_code=400, detail="Invalid file name")
    if Path(name).suffix.lower() not in RESOURCE_EXTS:
        raise HTTPException(status_code=400, detail="Only PDF, TXT, or Markdown files are allowed")
    return name


@app.get("/admin/resources")
def admin_list_resources(admin: dict = Depends(require_admin)):
    return {
        "files": _list_resource_files(),
        "index": {
            "rag_available": RAG_AVAILABLE,
            "total_chunks": len(_chunks),
            "sources": len({c.get("source") for c in _chunks}),
            "stale": _index_stale(),
        },
    }


@app.post("/admin/resources")
async def admin_upload_resource(file: UploadFile = File(...), admin: dict = Depends(require_admin)):
    name = _safe_resource_name(file.filename)
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    dest = DOCS_DIR / name
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")
    with open(dest, "wb") as f:
        f.write(data)
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "upload_resource", "", "", {"name": name, "size": len(data)}
    )
    return {"ok": True, "name": name, "size": len(data), "stale": _index_stale()}


@app.delete("/admin/resources/{filename}")
def admin_delete_resource(filename: str, admin: dict = Depends(require_admin)):
    name = _safe_resource_name(filename)
    target = DOCS_DIR / name
    if not target.exists():
        raise HTTPException(status_code=404, detail="File not found")
    target.unlink()
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "delete_resource", "", "", {"name": name}
    )
    return {"ok": True, "stale": _index_stale()}


@app.post("/admin/resources/rebuild")
def admin_rebuild_index(admin: dict = Depends(require_admin)):
    if not RAG_AVAILABLE:
        raise HTTPException(status_code=503, detail="Retrieval engine is not available on this server")
    stats = _build_index(force=True)
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "rebuild_knowledge_base", "", "", stats or {}
    )
    return {"ok": True, **(stats or {}), "stale": _index_stale()}


@app.get("/admin/resources/file/{filename}")
def admin_get_resource_file(filename: str, download: int = Query(0), admin: dict = Depends(require_admin)):
    """View (inline) or download a knowledge-base document."""
    from fastapi.responses import FileResponse
    name = _safe_resource_name(filename)
    target = DOCS_DIR / name
    if not target.exists():
        raise HTTPException(status_code=404, detail="File not found")
    media = {
        "pdf": "application/pdf",
        "txt": "text/plain",
        "md": "text/markdown",
        "markdown": "text/markdown",
    }.get(target.suffix.lower().lstrip("."), "application/octet-stream")
    disposition = "attachment" if download else "inline"
    return FileResponse(
        str(target), media_type=media, filename=name,
        headers={"Content-Disposition": f'{disposition}; filename="{name}"'},
    )


# ── Step resources (student Resources panel: video + interactive per step) ──

@app.get("/step-resources")
def public_step_resources(lang: Optional[str] = Query(None)):
    """Per-step video + interactive URLs, keyed by level, in the requested
    language. Blank non-English fields fall back to English. Read by students."""
    lang = lang if lang in SUPPORTED_LANGUAGES else "en"
    return {"resources": get_step_resources(lang), "lang": lang}


class StepResourceReq(BaseModel):
    step: int
    level: str
    lang: Optional[str] = "en"
    video_url: Optional[str] = ""
    interactive_url: Optional[str] = ""


@app.get("/admin/step-resources")
def admin_list_step_resources(admin: dict = Depends(require_admin)):
    """Raw stored values for every language (no fallback), for the editor."""
    return {"resources": get_step_resources_all(), "languages": sorted(SUPPORTED_LANGUAGES)}


@app.patch("/admin/step-resources")
def admin_update_step_resource(req: StepResourceReq, admin: dict = Depends(require_admin)):
    if req.level not in ("high_school", "higher_ed"):
        raise HTTPException(status_code=400, detail="Invalid education level")
    if not (1 <= req.step <= 9):
        raise HTTPException(status_code=400, detail="Step must be 1–9")
    lang = (req.lang or "en").strip().lower()
    if lang not in SUPPORTED_LANGUAGES:
        raise HTTPException(status_code=400, detail="Unsupported language")
    ok = upsert_step_resource(req.step, req.level, req.video_url or "",
                              req.interactive_url or "", lang)
    if not ok:
        raise HTTPException(status_code=400, detail="Could not save step resource")
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "update_step_resource", "", "", {"step": req.step, "level": req.level, "lang": lang}
    )
    return {"ok": True}


# ── Admin: Class Management ─────────────────────────────

@app.get("/admin/classes")
def admin_list_classes(
    skip: int = Query(0),
    limit: int = Query(50),
    search: str = Query(""),
    admin: dict = Depends(require_admin),
):
    classes, total = get_all_classes(skip, limit, search or None)
    # Sanitize ObjectId and hide password_hash
    for c in classes:
        c.pop("password_hash", None)
    return {"classes": classes, "total": total}


@app.delete("/admin/classes/{class_id}")
def admin_delete_class(class_id: str, admin: dict = Depends(require_admin)):
    from database import get_students_in_class as _get_students
    students = _get_students(class_id)
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "delete_class", class_id, "",
        {"student_count": len(students)}
    )
    ok = delete_class_by_id(class_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Class not found")
    return {"ok": True}


class AdminClassPatchReq(BaseModel):
    class_name: Optional[str] = None
    password: Optional[str] = None
    teacher_id: Optional[str] = None
    ai_enabled: Optional[bool] = None
    access_mode: Optional[str] = None


class AdminAddStudentsReq(BaseModel):
    count: int


class AdminCreateClassReq(BaseModel):
    teacher_id: str
    class_name: str
    password: str
    student_count: int


@app.get("/admin/classes/{class_id}")
def admin_get_class_detail(class_id: str, admin: dict = Depends(require_admin)):
    """Class drill-down: class info + per-student progress."""
    detail = get_class_detail(class_id)
    if not detail:
        raise HTTPException(status_code=404, detail="Class not found")
    return detail


@app.patch("/admin/classes/{class_id}")
def admin_update_class(
    class_id: str,
    req: AdminClassPatchReq,
    admin: dict = Depends(require_admin),
):
    """Admin class fixes: rename, change password (re-hashes all student
    logins), reassign teacher, override AI settings."""
    cls = find_class_by_id(class_id)
    if not cls:
        raise HTTPException(status_code=404, detail="Class not found")

    updates: dict = {}
    details: dict = {}
    if req.class_name is not None and req.class_name.strip():
        updates["class_name"] = req.class_name.strip()
        details["class_name"] = updates["class_name"]
    if req.password:
        if len(req.password) < 4:
            raise HTTPException(status_code=400, detail="Password must be at least 4 characters")
        pw_hash = hash_password(req.password)
        updates["password"] = req.password
        updates["password_hash"] = pw_hash
        details["students_rehashed"] = set_class_students_password(class_id, pw_hash)
        details["password_changed"] = True
    if req.teacher_id:
        target = find_user_by_id(req.teacher_id)
        if not target or target.get("role") != "teacher":
            raise HTTPException(status_code=400, detail="Target user is not a teacher")
        updates["teacher_id"] = req.teacher_id
        details["reassigned_to"] = target.get("email", "")

    settings_updates: dict = {}
    if req.ai_enabled is not None:
        settings_updates["ai_enabled"] = req.ai_enabled
    if req.access_mode is not None:
        if req.access_mode not in ("full", "step", "phase"):
            raise HTTPException(status_code=400, detail="access_mode must be full, step, or phase")
        settings_updates["access_mode"] = req.access_mode

    if not updates and not settings_updates:
        raise HTTPException(status_code=400, detail="Nothing to update")
    if updates:
        update_class_fields(class_id, updates)
    if settings_updates:
        update_class_settings(class_id, settings_updates)
        details["settings"] = settings_updates

    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "update_class", class_id, cls.get("class_code", ""), details)
    return {"ok": True}


@app.post("/admin/classes/{class_id}/students")
def admin_add_class_students(
    class_id: str,
    req: AdminAddStudentsReq,
    admin: dict = Depends(require_admin),
):
    """Add student slots to an existing class, continuing the numbering."""
    cls = find_class_by_id(class_id)
    if not cls:
        raise HTTPException(status_code=404, detail="Class not found")
    if not (1 <= req.count <= 100):
        raise HTTPException(status_code=400, detail="Count must be between 1 and 100")

    code = cls.get("class_code", "")
    max_n = 0
    for s in get_students_in_class(class_id):
        m = re.match(rf"^{re.escape(code)}_(\d+)$", s.get("username", ""))
        if m:
            max_n = max(max_n, int(m.group(1)))
    teacher = find_user_by_id(cls.get("teacher_id", ""))
    edu = (teacher or {}).get("education_level", "high_school")
    pw_hash = cls["password_hash"]
    created = []
    for i in range(max_n + 1, max_n + 1 + req.count):
        username = f"{code}_{i:02d}"
        create_classroom_student(username, pw_hash, f"Student {i:02d}", class_id, education_level=edu)
        created.append({"username": username, "name": f"Student {i:02d}"})
    update_class_fields(class_id, {"student_count": int(cls.get("student_count", 0)) + req.count})

    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "add_class_students", class_id, code, {"count": req.count})
    return {"ok": True, "students": created}


@app.post("/admin/classes")
def admin_create_class(req: AdminCreateClassReq, admin: dict = Depends(require_admin)):
    """Create a class on behalf of a teacher (same flow as /teacher/create-class)."""
    teacher = find_user_by_id(req.teacher_id)
    if not teacher or teacher.get("role") != "teacher":
        raise HTTPException(status_code=400, detail="Target user is not a teacher")
    if not (1 <= req.student_count <= 100):
        raise HTTPException(status_code=400, detail="Student count must be between 1 and 100")
    if len(req.password) < 4:
        raise HTTPException(status_code=400, detail="Password must be at least 4 characters")

    raw = re.sub(r'[^a-z0-9]', '', req.class_name.lower().replace(' ', ''))
    class_code = raw[:20] or "class"
    base_code, counter = class_code, 1
    while find_class_by_code(class_code):
        class_code = f"{base_code}{counter}"
        counter += 1

    pw_hash = hash_password(req.password)
    class_id = create_class_doc(
        str(teacher["_id"]), req.class_name, class_code, pw_hash, req.password, req.student_count)
    edu = teacher.get("education_level", "high_school")
    for i in range(1, req.student_count + 1):
        create_classroom_student(
            f"{class_code}_{i:02d}", pw_hash, f"Student {i:02d}", class_id, education_level=edu)

    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "create_class", class_id, class_code,
        {"teacher": teacher.get("email", ""), "students": req.student_count})
    return {"ok": True, "class_id": class_id, "class_code": class_code}


# ── Admin: Session / Design Viewer ──────────────────────

@app.get("/admin/sessions")
def admin_list_sessions(
    skip: int = Query(0),
    limit: int = Query(50),
    user_id: str = Query(""),
    search: str = Query(""),
    step: int = Query(0),
    status: str = Query(""),
    sort_by: str = Query("updated_at"),
    sort_dir: str = Query("desc"),
    admin: dict = Depends(require_admin),
):
    sessions, total = get_all_sessions(
        skip, limit, user_id or None, search or None,
        step or None, status or None, sort_by, sort_dir)
    return {"sessions": sessions, "total": total, "stats": get_session_stats()}


@app.delete("/admin/sessions/{session_id}")
def admin_delete_session(session_id: str, admin: dict = Depends(require_admin)):
    sess = find_session(session_id)
    if not sess:
        raise HTTPException(status_code=404, detail="Session not found")
    owner = find_user_by_id(sess.get("user_id", ""))
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "delete_session", session_id,
        (owner or {}).get("email") or (owner or {}).get("username", ""),
        {"active_step": sess.get("active_step")})
    delete_session_by_sid(session_id)
    return {"ok": True}


class SessionCleanupReq(BaseModel):
    mode: str  # "empty" | "orphaned"


@app.post("/admin/sessions/cleanup")
def admin_cleanup_sessions(req: SessionCleanupReq, admin: dict = Depends(require_admin)):
    """Bulk-delete empty (no step work) or orphaned (owner deleted) sessions."""
    if req.mode not in ("empty", "orphaned"):
        raise HTTPException(status_code=400, detail="mode must be 'empty' or 'orphaned'")
    deleted = delete_sessions_bulk(req.mode)
    record_admin_action(
        str(admin["_id"]), admin.get("email", ""),
        "cleanup_sessions", req.mode, "", {"deleted": deleted})
    return {"ok": True, "deleted": deleted}


@app.get("/admin/sessions/{session_id}")
def admin_get_session(session_id: str, admin: dict = Depends(require_admin)):
    """Get full session data for admin viewer (reuses teacher endpoint logic)."""
    raw = find_session(session_id)
    if not raw:
        raise HTTPException(status_code=404, detail="Session not found")
    owner = find_user_by_id(raw.get("user_id", ""))
    return {
        "session_id": raw.get("session_id"),
        "user_name": (owner.get("username") or owner.get("name", "")) if owner else "",
        "user_email": (owner.get("email") or owner.get("username", "")) if owner else "",
        "active_step": raw.get("active_step"),
        "worldview_label": raw.get("worldview_label"),
        "resolved_path": raw.get("resolved_path"),
        "chosen_methodology": raw.get("chosen_methodology"),
        "step_notes": raw.get("step_notes", {}),
        "created_at": raw.get("created_at"),
        "updated_at": raw.get("updated_at"),
    }


# ── Admin: System Health ────────────────────────────────

@app.get("/admin/health")
def admin_system_health(admin: dict = Depends(require_admin)):
    import time
    health = {
        "server": "ok",
        "uptime_seconds": int(time.time() - _SERVER_START_TIME),
    }

    # MongoDB
    try:
        from database import client as mongo_client
        mongo_client.admin.command("ping")
        health["mongodb"] = "ok"
    except Exception as e:
        health["mongodb"] = f"error: {e}"

    # LLM backends
    health["llm_backend"] = LLM_BACKEND

    # vLLM health
    try:
        vllm_health_url = VLLM_URL.replace("/v1/chat/completions", "/health")
        r = requests.get(vllm_health_url, timeout=3)
        health["vllm"] = "ok" if r.status_code == 200 else f"status {r.status_code}"
    except Exception as e:
        health["vllm"] = f"error: {e}"

    # Ollama health
    try:
        r = requests.get(f"{OLLAMA_BASE}/api/tags", timeout=3)
        models = [m.get("name", "") for m in r.json().get("models", [])]
        health["ollama"] = "ok"
        health["ollama_models"] = models
    except Exception as e:
        health["ollama"] = f"error: {e}"
        health["ollama_models"] = []

    # RAG
    health["rag_available"] = RAG_AVAILABLE
    health["rag_index_loaded"] = _faiss_index is not None

    # Disk
    import shutil
    usage = shutil.disk_usage("/")
    health["disk_total_gb"] = round(usage.total / (1024**3), 1)
    health["disk_free_gb"] = round(usage.free / (1024**3), 1)

    # CPU / RAM / load
    try:
        import psutil
        health["cpu_percent"] = psutil.cpu_percent(interval=0.2)
        health["cpu_count"] = psutil.cpu_count()
        mem = psutil.virtual_memory()
        health["ram_total_gb"] = round(mem.total / (1024**3), 1)
        health["ram_used_gb"] = round(mem.used / (1024**3), 1)
        health["ram_percent"] = mem.percent
        health["load_avg"] = [round(x, 2) for x in os.getloadavg()]
    except Exception:
        pass

    # GPUs (empty list when nvidia-smi is unavailable/broken - itself a signal,
    # since a stale driver silently pushes Ollama onto CPU)
    health["gpus"] = []
    try:
        import subprocess
        out = subprocess.run(
            ["nvidia-smi", "--query-gpu=name,memory.used,memory.total,utilization.gpu,temperature.gpu",
             "--format=csv,noheader,nounits"],
            capture_output=True, text=True, timeout=5)
        if out.returncode == 0:
            for line in out.stdout.strip().splitlines():
                parts = [p.strip() for p in line.split(",")]
                if len(parts) == 5:
                    health["gpus"].append({
                        "name": parts[0],
                        "mem_used_mb": int(float(parts[1])),
                        "mem_total_mb": int(float(parts[2])),
                        "util_percent": int(float(parts[3])),
                        "temp_c": int(float(parts[4])),
                    })
    except Exception:
        pass

    # Data volumes
    try:
        health["db_users"] = sum(get_user_counts_by_role().values())
        health["db_sessions"] = get_total_sessions_count()
    except Exception:
        pass

    return health


@app.post("/admin/health/llm-test")
def admin_llm_latency_test(admin: dict = Depends(require_admin)):
    """Round-trip a tiny prompt through the active LLM backend and time it."""
    import time as _time
    start = _time.time()
    try:
        r = requests.post(f"{OLLAMA_BASE}/api/generate", json={
            "model": LLM_MODEL,
            "prompt": "Reply with the single word: ok",
            "stream": False,
            "options": {"num_predict": 5},
        }, timeout=60)
        r.raise_for_status()
        latency = round(_time.time() - start, 2)
        return {"ok": True, "latency_seconds": latency, "model": LLM_MODEL,
                "reply": (r.json().get("response") or "").strip()[:40]}
    except Exception as e:
        return {"ok": False, "latency_seconds": round(_time.time() - start, 2),
                "model": LLM_MODEL, "error": str(e)[:200]}


# ── Admin: CSV Data Export ──────────────────────────────

@app.get("/admin/export/users.csv")
def admin_export_users_csv(admin: dict = Depends(require_admin)):
    import csv, io
    users, _ = get_all_users(skip=0, limit=100000)
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["Name", "Email", "Username", "Role", "Education", "Active", "Created", "Last Login"])
    for u in users:
        writer.writerow([
            u.get("name", ""), u.get("email", ""), u.get("username", ""),
            u.get("role", ""), u.get("education_level", ""),
            "Yes" if u.get("is_active", True) else "No",
            u.get("created_at", ""), u.get("last_login_at", ""),
        ])
    return Response(
        content=buf.getvalue(),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=hopscotch_users.csv"},
    )


@app.get("/admin/export/sessions.csv")
def admin_export_sessions_csv(admin: dict = Depends(require_admin)):
    import csv, io
    sessions, _ = get_all_sessions(skip=0, limit=100000)
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["Session ID", "User", "Email", "Active Step", "Worldview", "Path", "Methodology", "Created", "Updated"])
    for s in sessions:
        writer.writerow([
            s.get("session_id", ""), s.get("user_name", ""), s.get("user_email", ""),
            s.get("active_step", ""), s.get("worldview_label", ""),
            s.get("resolved_path", ""), s.get("chosen_methodology", ""),
            s.get("created_at", ""), s.get("updated_at", ""),
        ])
    return Response(
        content=buf.getvalue(),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=hopscotch_sessions.csv"},
    )


@app.get("/admin/export/logins.csv")
def admin_export_logins_csv(admin: dict = Depends(require_admin)):
    import csv, io
    logins = get_recent_logins(limit=100000)
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["User", "IP", "City", "Region", "Country", "Time", "Success"])
    for l in logins:
        writer.writerow([
            l.get("email", ""), l.get("ip", ""),
            l.get("city", ""), l.get("region", ""), l.get("country", ""),
            l.get("login_at", ""), "Yes" if l.get("success") else "No",
        ])
    return Response(
        content=buf.getvalue(),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=hopscotch_logins.csv"},
    )


# ── Admin: User Detail Drill-Down ──────────────────────

@app.get("/admin/users/{user_id}/detail")
def admin_user_detail(user_id: str, admin: dict = Depends(require_admin)):
    detail = get_user_detail(user_id)
    if not detail:
        raise HTTPException(status_code=404, detail="User not found")
    return detail


def _replace_pptx_text(shape, replacements: dict):
    """Recursively replace placeholder text in a PPTX shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Handle group shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _replace_pptx_text(child, replacements)
        return

    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for placeholder, value in replacements.items():
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)
